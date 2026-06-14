from __future__ import annotations

import argparse
import json
import math
import os
from pathlib import Path
import re
import shutil
from typing import Any, Iterable

from PIL import Image, ImageDraw, ImageFilter, ImageFont, ImageOps
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[4]
WORK_DIR = Path(__file__).resolve().parent
ASSET_DIR = ROOT / "docs" / "数据库系统基础" / "最终交付" / "07_Gemini交付包" / "02_Gemini_生成豆包PPT包" / "上传资料" / "04_界面截图与图表素材"
OUTPUT_DIR = WORK_DIR / "page_visuals"
BG_DIR = WORK_DIR / "page_backgrounds"
PROMPT_PATH = WORK_DIR / "image_generation_prompts.jsonl"
REVIEW_HTML = WORK_DIR / "review_rescue_deck.html"
PPTX_RESCUE = WORK_DIR / "第07组毕振岚-数据库课程设计答辩PPT_融合救场版.pptx"
PPTX_FINAL = ROOT / "docs" / "数据库系统基础" / "最终交付" / "03_答辩PPT" / "第07组毕振岚-数据库课程设计答辩PPT.pptx"
REVIEW_SHELL = Path(r"D:\Documents\Codex-Migrated\skills\ppt-image-first\assets\review_shell\index.html")

W, H = 1600, 900
SLIDE_W, SLIDE_H = 13.333333, 7.5

INK = (13, 27, 54)
BLUE = (28, 115, 255)
CYAN = (51, 183, 214)
GREEN = (84, 216, 166)
MUTED = (96, 113, 137)
ORANGE = (255, 140, 58)
PANEL = (255, 255, 255, 218)
PANEL_SOFT = (255, 255, 255, 188)
LINE = (208, 222, 238, 220)


def font_path(*candidates: str) -> str:
    for raw in candidates:
        path = Path(raw)
        if path.exists():
            return str(path)
    return "arial.ttf"


FONT_REG = font_path(r"C:\Windows\Fonts\msyh.ttc", r"C:\Windows\Fonts\simhei.ttf")
FONT_BOLD = font_path(r"C:\Windows\Fonts\msyhbd.ttc", r"C:\Windows\Fonts\msyh.ttc")
FONT_SEGOE = font_path(r"C:\Windows\Fonts\seguisb.ttf", r"C:\Windows\Fonts\segoeui.ttf", r"C:\Windows\Fonts\msyh.ttc")


def f(size: int, bold: bool = False, latin: bool = False) -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(FONT_SEGOE if latin else (FONT_BOLD if bold else FONT_REG), size)


SLIDES: list[dict[str, Any]] = [
    {
        "id": "S01",
        "title": "基于 PostgreSQL 的数字图像集成管理系统",
        "subtitle": "JavaFX + PostgreSQL 桌面图片管理｜数据库系统课程设计答辩",
        "role": "封面",
        "visual": "cover",
        "metrics": [("13", "物理表"), ("19", "索引"), ("4", "视图"), ("5", "触发器"), ("5", "存储过程")],
        "assets": ["real_02_主界面_1200x800_默认窗口.png"],
    },
    {
        "id": "S02",
        "title": "答辩先对齐得分点",
        "subtitle": "数据库设计是主线，功能、后台和界面都服务于落库与查询",
        "role": "评分点总览",
        "visual": "score",
    },
    {
        "id": "S03",
        "title": "系统功能结构：图片管理不是文件夹外壳",
        "subtitle": "目录、图片、标签、AI、版本、日志共同组成可查询的数据资产",
        "role": "功能结构",
        "visual": "asset_focus",
        "assets": ["图1_系统功能结构图.png", "real_02_主界面_1200x800_默认窗口.png"],
    },
    {
        "id": "S04",
        "title": "ER 模型：围绕 images 主表展开",
        "subtitle": "目录自引用、标签多对多、版本一对多、操作日志可追踪",
        "role": "数据库设计",
        "visual": "asset_focus",
        "assets": ["图2_ER关系模式图.png"],
    },
    {
        "id": "S05",
        "title": "物理结构：13 张表按职责分层",
        "subtitle": "基础对象、AI 标签、历史审计、扩展预留边界清晰",
        "role": "表结构",
        "visual": "tables",
    },
    {
        "id": "S06",
        "title": "高级数据库对象全景",
        "subtitle": "索引、视图、触发器、存储过程和递归 CTE 共同支撑性能与一致性",
        "role": "SQL 对象",
        "visual": "asset_focus",
        "assets": ["图4_SQL对象与性能设计.png"],
    },
    {
        "id": "S07",
        "title": "完整磁盘目录树：懒加载 + 递归 CTE",
        "subtitle": "前台只展开当前需要的节点，数据库保留完整目录层级与路径统计",
        "role": "目录树",
        "visual": "directory",
        "assets": ["real_03_主界面_1440x900_宽屏窗口.png"],
    },
    {
        "id": "S08",
        "title": "图片入库链路：缩略图、元数据与 bytea",
        "subtitle": "扫描目录后把文件系统信息结构化，缩略图缓存直接服务主界面渲染",
        "role": "持久化",
        "visual": "pipeline",
        "assets": ["real_09_图片查看器_960x680.png"],
    },
    {
        "id": "S09",
        "title": "后台基座：连接池、预编译与事务回滚",
        "subtitle": "HikariCP 管连接，PreparedStatement 防注入，批量操作失败即回滚",
        "role": "后台程序设计",
        "visual": "asset_focus",
        "assets": ["图6_后台架构图.png"],
    },
    {
        "id": "S10",
        "title": "AI 标签与 NL2SQL：能问，但只能安全地问",
        "subtitle": "AI 结果先落库；自然语言查询只允许走 v_image_search 只读视图",
        "role": "AI 搜索",
        "visual": "asset_focus",
        "assets": ["图5_NL2SQL安全链路.png", "05_标签与扩展搜索.png"],
    },
    {
        "id": "S11",
        "title": "版本历史：从编辑行为回到数据库一致性",
        "subtitle": "image_versions 记录快照，sp_restore_version 支撑恢复，operation_logs 留痕",
        "role": "版本与审计",
        "visual": "version",
        "assets": ["real_14_图片编辑器_1250x938.png", "04_图片编辑与版本历史.png"],
    },
    {
        "id": "S12",
        "title": "界面证据：核心流程都已能演示",
        "subtitle": "启动向导、主界面、图片查看、幻灯片、编辑器、数据库初始化",
        "role": "界面设计",
        "visual": "gallery",
        "assets": [
            "real_04_首次启动向导_640x620.png",
            "real_08_数据库连接与初始化向导_760x680.png",
            "real_10_图片查看器_1200x850.png",
            "real_12_幻灯片播放_1250x875.png",
            "real_14_图片编辑器_1250x938.png",
            "real_16_批量重命名_563x450.png",
        ],
    },
    {
        "id": "S13",
        "title": "现场演示路线",
        "subtitle": "先证明数据库初始化，再展示落库、查询、版本和幻灯片主流程",
        "role": "演示路线",
        "visual": "demo",
        "assets": ["图7_演示路线.png"],
    },
    {
        "id": "S14",
        "title": "总结：用数据库重新组织本地图片",
        "subtitle": "已完成主流程，云端/WebDAV、语音搜索作为后续扩展预留",
        "role": "总结",
        "visual": "summary",
    },
]


def ensure_dirs() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    BG_DIR.mkdir(parents=True, exist_ok=True)


def write_planning_files() -> None:
    content_report = """# content_report.md

本答辩 PPT 的内容基础来自课程任务书、数据库课程设计报告、README、schema.sql、真实界面截图和最终交付清单。材料已经足够支撑完整答辩，因此本阶段采取“整理但不扩写”的方式：保留项目真实实现，把表达重心压到数据库课程评分点上。

这套系统的核心论点不是“做了一个图片查看器”，而是“用 PostgreSQL 把本地图片目录、图片元数据、缩略图、标签、AI 分析结果、搜索历史、版本历史和操作日志统一组织成可维护、可查询、可追踪的数据资产”。JavaFX 负责交互，PostgreSQL 负责结构化、约束、索引、视图、触发器、存储过程和事务一致性。答辩需要先让老师看到数据库设计的工作量，再说明功能与界面如何从这些数据库对象中受益。

可视化内容应优先使用三类证据：一是 schema.sql 中明确存在的表、视图、索引、触发器和存储过程；二是真实运行截图，如主界面、数据库初始化向导、图片编辑器和幻灯片；三是已经生成的 ER 图、后台架构图、NL2SQL 安全链路图和演示路线图。云端/WebDAV、语音搜索、SQLite 多底层切换只能作为扩展预留，不进入已实现主流程。
"""
    design_spec = """# design_spec.md

项目名称：基于 PostgreSQL 的数字图像集成管理系统。
使用场景：数据库系统课程设计 15 分钟答辩。
受众：数据库课程老师与同学评委。
页面范围：14 页，16:9。

叙事主线：先以评分点总览建立“数据库含金量”预期，再展示 ER 模型、物理表结构和高级 SQL 对象，随后解释 JavaFX 后台如何通过 HikariCP、PreparedStatement 和事务控制安全访问数据库，最后用真实截图证明功能可演示，并诚实说明扩展边界。

视觉方向：参考同目录“收齐其他小组作业/PPT.pptx”的 image-first 答辩风格。整套 PPT 保持亮色、冰蓝、深墨蓝、青绿色强调线；背景使用整页生成式 UI/数据库场景，不使用廉价深色科技风。正文页采用大标题、短句、真实截图或数据库图表，避免密集文字。

deck-level continuity anchor：
- brightness_world: high-key bright product-defense board
- background_tendency: white / ice-blue / soft 3D database and desktop-app scene
- palette_roles: deep navy for title, vivid blue for structure, cyan-green for proof accent, orange only for warnings and boundaries
- lighting_model: soft luminous depth, broad shadows, no neon cyberpunk
- material_language: translucent white panels, rounded app windows, database cylinders, connected relation cards
- text_policy: large accurate Chinese text is rendered by deterministic slide layer; generated backgrounds must avoid random small text
"""
    blueprint_lines = ["# slide_blueprint.md", ""]
    for slide in SLIDES:
        blueprint_lines.append(f"## {slide['id']} {slide['title']}")
        blueprint_lines.append(f"- page_role: {slide['role']}")
        blueprint_lines.append(f"- core_message: {slide['subtitle']}")
        blueprint_lines.append("- content_basis_binding: README、schema_object_summary.json、真实截图和课程评分点。")
        blueprint_lines.append("- claim_status: user_provided / repository_verified。")
        blueprint_lines.append("- visual_strategy: 整页生成式背景承载氛围，准确中文和数据库对象由最终页面图层呈现。")
        blueprint_lines.append("- continuity_inheritance: 继承亮色冰蓝背景、深墨蓝标题、蓝绿强调线、圆角应用窗口和数据库卡片语法。")
        blueprint_lines.append("")
    spec_lock = """# spec_lock.md

canvas:
- format: PPTX
- ratio: 16:9
- rendered_page_size: 1600x900

visual_system:
- chosen_style_direction: bright image-first database defense deck
- no dark hacker background
- no dense Doubao-style full paragraphs
- no random private chat, prompt, screenshot-source wording in final slides
- generated background may not contain random tiny pseudo text
- final accurate text must come from approved slide blueprint

content_grounding:
- database object counts: 13 tables, 19 indexes, 4 views, 4 functions, 5 triggers, 5 procedures
- cloud/WebDAV only extension reserve
- AI endpoint requires user configuration and supports offline fallback
- NL2SQL limited to read-only SELECT query path over v_image_search
"""
    (WORK_DIR / "content_report.md").write_text(content_report, encoding="utf-8")
    (WORK_DIR / "design_spec.md").write_text(design_spec, encoding="utf-8")
    (WORK_DIR / "slide_blueprint.md").write_text("\n".join(blueprint_lines), encoding="utf-8")
    (WORK_DIR / "spec_lock.md").write_text(spec_lock, encoding="utf-8")


def background_prompt(slide: dict[str, Any]) -> str:
    title = slide["title"]
    role = slide["role"]
    scene = {
        "cover": "hero cover with a PostgreSQL database cylinder, elegant desktop photo-management application mockup, connected metadata cards, spatial depth",
        "score": "scorecard overview with database core in the center, six assessment modules orbiting as polished panels",
        "asset_focus": "large clean presentation canvas with one dominant app-window or database-diagram area, subtle connected cards and data lines",
        "tables": "database schema wall with grouped table cards, relation lines, PostgreSQL-like data storage motif",
        "directory": "file-system tree and recursive hierarchy visual, lazy-loading nodes flowing into PostgreSQL database storage",
        "pipeline": "image ingestion pipeline from folder to thumbnail bytea cache to metadata table to UI grid",
        "version": "photo editor history timeline, snapshots, restore arrow and audit log cards",
        "gallery": "polished multi-window desktop application evidence wall, six screenshot frames with depth",
        "demo": "demo route map through initialization, scanning, search, version restore and slideshow",
        "summary": "final summary scene with stable database core, desktop app, extension modules fading into future reserve",
    }.get(slide["visual"], "database course defense presentation scene")
    return (
        "Use case: productivity-visual\n"
        "Asset type: full 16:9 slide background for a Chinese database course defense PPT\n"
        f"Primary request: Create a rich image-first presentation background for the slide role '{role}' and topic '{title}'.\n"
        f"Scene/backdrop: {scene}.\n"
        "Style/medium: bright high-key image-first PPT, like a polished generated presentation page, white and ice-blue academic product style, soft 3D depth, translucent rounded application windows, database cylinders, relation cards, fine blue connector lines, subtle green accents.\n"
        "Composition/framing: leave clear readable zones for large title and exact overlay text; keep the page visually full and premium, not empty; no dark cyberpunk style.\n"
        "Lighting/mood: soft daylight, clean academic defense, professional and expensive-looking, no harsh neon.\n"
        "Color palette: white, pale ice-blue, deep navy accents, vivid blue connectors, small cyan-green highlights.\n"
        "Constraints: no logos except generic PostgreSQL elephant-like database motif if needed, no watermark, no random tiny text, no gibberish labels, no English lorem ipsum, no QR code, no people.\n"
        "Avoid: dark hacker interface, purple gradient template, cheap stock illustration, crowded paragraph text, excessive icons."
    )


def write_prompts() -> None:
    ensure_dirs()
    rows = []
    for slide in SLIDES:
        rows.append({
            "prompt": background_prompt(slide),
            "out": f"{slide['id']}.png",
            "size": "2048x1152",
            "quality": "medium",
            "output_format": "png",
        })
    PROMPT_PATH.write_text("\n".join(json.dumps(row, ensure_ascii=False) for row in rows) + "\n", encoding="utf-8")


def crop_cover(img: Image.Image, size: tuple[int, int] = (W, H)) -> Image.Image:
    return ImageOps.fit(img.convert("RGB"), size, method=Image.Resampling.LANCZOS, centering=(0.5, 0.5))


def load_bg(slide_id: str) -> Image.Image:
    path = BG_DIR / f"{slide_id}.png"
    if path.exists():
        return crop_cover(Image.open(path))
    img = Image.new("RGB", (W, H), (241, 248, 255))
    d = ImageDraw.Draw(img, "RGBA")
    d.ellipse((-210, -250, 420, 260), fill=(214, 232, 250, 230))
    d.ellipse((1160, -190, 1790, 260), fill=(223, 239, 255, 210))
    d.ellipse((-180, 710, 460, 1090), fill=(92, 216, 190, 120))
    d.ellipse((215, 780, 820, 1130), fill=(28, 115, 255, 105))
    return img


def add_wash(img: Image.Image, opacity: int = 48) -> Image.Image:
    layer = Image.new("RGBA", img.size, (255, 255, 255, opacity))
    return Image.alpha_composite(img.convert("RGBA"), layer)


def shadowed_round(draw_img: Image.Image, box: tuple[int, int, int, int], radius: int, fill: tuple[int, int, int, int] = PANEL, outline: tuple[int, int, int, int] | None = LINE, shadow: int = 18) -> None:
    x1, y1, x2, y2 = box
    if shadow:
        sh = Image.new("RGBA", draw_img.size, (0, 0, 0, 0))
        sd = ImageDraw.Draw(sh, "RGBA")
        sd.rounded_rectangle((x1, y1 + 8, x2, y2 + 8), radius=radius, fill=(31, 78, 130, 38))
        sh = sh.filter(ImageFilter.GaussianBlur(shadow))
        draw_img.alpha_composite(sh)
    d = ImageDraw.Draw(draw_img, "RGBA")
    d.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=1)


def draw_text(draw: ImageDraw.ImageDraw, xy: tuple[int, int], text: str, size: int, color=INK, bold=False, max_width: int | None = None, line_gap: int = 8, latin: bool = False) -> int:
    font = f(size, bold, latin=latin)
    if not max_width:
        draw.text(xy, text, font=font, fill=color)
        return int(draw.textbbox(xy, text, font=font)[3])
    lines: list[str] = []
    cur = ""
    for ch in text:
        nxt = cur + ch
        if draw.textlength(nxt, font=font) <= max_width or not cur:
            cur = nxt
        else:
            lines.append(cur)
            cur = ch
    if cur:
        lines.append(cur)
    y = xy[1]
    for line in lines:
        draw.text((xy[0], y), line, font=font, fill=color)
        y += size + line_gap
    return y


def title_block(img: Image.Image, title: str, subtitle: str, top: int = 42, width: int = 1220) -> None:
    d = ImageDraw.Draw(img, "RGBA")
    shadowed_round(img, (36, 24, min(1538, width + 128), 160), 26, fill=(255, 255, 255, 202), shadow=8)
    draw_text(d, (54, top), title, 44, INK, True, max_width=width, line_gap=6)
    d.rounded_rectangle((56, top + 72, 204, top + 80), radius=4, fill=BLUE)
    d.rounded_rectangle((204, top + 72, 258, top + 80), radius=4, fill=GREEN)
    draw_text(d, (56, top + 96), subtitle, 20, MUTED, False, max_width=width)


def paste_asset(img: Image.Image, name: str, box: tuple[int, int, int, int], radius: int = 22, shadow: int = 20, fit: str = "contain") -> None:
    path = ASSET_DIR / name
    if not path.exists():
        return
    x1, y1, x2, y2 = box
    target = (x2 - x1, y2 - y1)
    src = Image.open(path).convert("RGB")
    if fit == "cover":
        src = ImageOps.fit(src, target, Image.Resampling.LANCZOS)
    else:
        src.thumbnail(target, Image.Resampling.LANCZOS)
        canvas = Image.new("RGB", target, (250, 253, 255))
        canvas.paste(src, ((target[0] - src.width) // 2, (target[1] - src.height) // 2))
        src = canvas
    mask = Image.new("L", target, 0)
    md = ImageDraw.Draw(mask)
    md.rounded_rectangle((0, 0, target[0], target[1]), radius=radius, fill=255)
    shadowed_round(img, box, radius=radius, fill=(255, 255, 255, 210), shadow=shadow)
    img.paste(src, (x1, y1), mask)


def metric_card(img: Image.Image, box: tuple[int, int, int, int], value: str, label: str, accent=BLUE) -> None:
    shadowed_round(img, box, 24, fill=(255, 255, 255, 220), shadow=16)
    d = ImageDraw.Draw(img, "RGBA")
    x1, y1, x2, y2 = box
    d.ellipse((x1 + 18, y1 + 20, x1 + 58, y1 + 60), fill=(*accent, 36) if len(accent) == 3 else accent)
    draw_text(d, (x1 + 76, y1 + 20), value, 34, INK, True)
    draw_text(d, (x1 + 76, y1 + 62), label, 17, MUTED)


def bullet_panel(img: Image.Image, box: tuple[int, int, int, int], heading: str, bullets: Iterable[str], accent=BLUE) -> None:
    shadowed_round(img, box, 26, fill=PANEL, shadow=18)
    d = ImageDraw.Draw(img, "RGBA")
    x1, y1, x2, y2 = box
    d.rounded_rectangle((x1 + 22, y1 + 24, x1 + 30, y1 + 62), radius=4, fill=accent)
    draw_text(d, (x1 + 44, y1 + 22), heading, 26, INK, True, max_width=x2 - x1 - 72)
    y = y1 + 78
    for item in bullets:
        d.ellipse((x1 + 44, y + 10, x1 + 54, y + 20), fill=accent)
        y = draw_text(d, (x1 + 70, y), item, 19, (38, 54, 78), False, max_width=x2 - x1 - 104, line_gap=5) + 14


def draw_cover(img: Image.Image, slide: dict[str, Any]) -> None:
    d = ImageDraw.Draw(img, "RGBA")
    shadowed_round(img, (64, 82, 672, 658), 34, fill=(255, 255, 255, 184), shadow=10)
    draw_text(d, (104, 126), "数据库系统课程设计答辩", 24, MUTED)
    draw_text(d, (104, 200), "基于 PostgreSQL 的", 52, INK, True, max_width=520)
    draw_text(d, (104, 276), "数字图像集成", 52, INK, True, max_width=520)
    draw_text(d, (104, 352), "管理系统", 52, INK, True, max_width=520)
    d.rounded_rectangle((108, 456, 266, 466), radius=5, fill=BLUE)
    d.rounded_rectangle((266, 456, 330, 466), radius=5, fill=GREEN)
    draw_text(d, (108, 502), "JavaFX 桌面端 · PostgreSQL 数据库驱动", 24, (42, 65, 98), max_width=470)
    draw_text(d, (108, 564), "第07组  毕振岚 · 陈厚华 · 徐阳", 24, INK, True, max_width=470)
    for i, (value, label) in enumerate(slide["metrics"]):
        metric_card(img, (690 + (i % 3) * 258, 600 + (i // 3) * 112, 920 + (i % 3) * 258, 690 + (i // 3) * 112), value, label)
    paste_asset(img, "real_02_主界面_1200x800_默认窗口.png", (710, 122, 1508, 560), radius=30, fit="cover")
    draw_text(d, (64, 806), "数据 · 图像 · 标签 · 搜索 · 版本 · 安全", 22, (42, 82, 133), True)


def draw_score(img: Image.Image) -> None:
    title_block(img, "答辩先对齐得分点", "把数据库设计放在视觉和讲述的中心位置")
    items = [
        ("30", "数据库设计", "ER、关系模式、外模式、索引、视图、触发器、过程"),
        ("20", "功能设计", "目录树、缩略图、标签、搜索、版本、幻灯片"),
        ("15", "后台程序", "HikariCP、PreparedStatement、事务回滚、离线降级"),
        ("15", "界面设计", "真实 JavaFX 截图证明主流程可用"),
        ("10", "报告表述", "按模板完整解释设计依据与测试"),
        ("10", "PPT表达", "少字、图优先、直入数据库含金量"),
    ]
    positions = [(78, 190), (565, 190), (1052, 190), (78, 500), (565, 500), (1052, 500)]
    for (score, name, desc), (x, y) in zip(items, positions):
        shadowed_round(img, (x, y, x + 420, y + 222), 32, fill=(255, 255, 255, 218), shadow=20)
        d = ImageDraw.Draw(img, "RGBA")
        d.ellipse((x + 28, y + 34, x + 116, y + 122), fill=(28, 115, 255, 28), outline=(28, 115, 255, 150), width=2)
        draw_text(d, (x + 49, y + 48), score, 34, BLUE, True, latin=True)
        draw_text(d, (x + 142, y + 40), name, 30, INK, True)
        draw_text(d, (x + 142, y + 92), desc, 19, MUTED, max_width=238, line_gap=6)
        d.rounded_rectangle((x + 142, y + 170, x + 310, y + 178), radius=4, fill=BLUE if score != "15" else GREEN)


def draw_asset_focus(img: Image.Image, slide: dict[str, Any]) -> None:
    if slide["id"] == "S06":
        draw_sql_objects(img)
        return
    title_block(img, slide["title"], slide["subtitle"])
    assets = slide.get("assets", [])
    if len(assets) == 1:
        paste_asset(img, assets[0], (92, 182, 1508, 814), radius=30, fit="contain")
    elif len(assets) >= 2:
        paste_asset(img, assets[0], (74, 188, 940, 810), radius=30, fit="contain")
        paste_asset(img, assets[1], (990, 228, 1512, 760), radius=30, fit="cover")
    if slide["id"] == "S10":
        bullet_panel(img, (1006, 628, 1490, 816), "安全边界", ["只读视图 v_image_search", "仅允许 SELECT", "超时与最大行数限制"], accent=GREEN)
    if slide["id"] == "S09":
        bullet_panel(img, (1038, 612, 1490, 816), "后台防线", ["连接池复用连接", "预编译 SQL 防注入", "批量失败主动 rollback"], accent=BLUE)


def draw_tables(img: Image.Image) -> None:
    title_block(img, "物理结构：13 张表按职责分层", "表结构不是堆名字，而是把图像管理拆成四个稳定域")
    groups = [
        ("基础目录与图片", ["directories", "images", "app_settings"], BLUE),
        ("标签与 AI 结果", ["tag_categories", "tags", "image_tags", "ai_analysis_results"], GREEN),
        ("版本与操作审计", ["image_versions", "image_edit_operations", "operation_logs", "search_history"], CYAN),
        ("扩展预留边界", ["cloud_sources", "cloud_images"], ORANGE),
    ]
    x_positions = [72, 440, 808, 1176]
    for (name, tables, color), x in zip(groups, x_positions):
        shadowed_round(img, (x, 198, x + 322, 744), 30, fill=(255, 255, 255, 216), shadow=18)
        d = ImageDraw.Draw(img, "RGBA")
        d.rounded_rectangle((x + 28, 226, x + 126, 236), radius=4, fill=color)
        draw_text(d, (x + 28, 258), name, 28, INK, True, max_width=264)
        y = 328
        for table in tables:
            d.rounded_rectangle((x + 28, y, x + 294, y + 58), radius=16, fill=(245, 250, 255, 235), outline=(207, 224, 242, 220))
            draw_text(d, (x + 48, y + 15), table, 20, (35, 78, 138), True, latin=True)
            y += 76
    d = ImageDraw.Draw(img, "RGBA")
    draw_text(d, (88, 790), "主流程已完成：本地目录、图片元数据、缩略图、标签、AI结果、版本历史、日志均可落库", 24, INK, True, max_width=1290)
    draw_text(d, (88, 830), "扩展预留：cloud_sources / cloud_images 只作为 WebDAV 与云端同步后续接口，不在答辩中夸大为已完成主流程", 20, MUTED, max_width=1320)


def draw_sql_objects(img: Image.Image) -> None:
    title_block(img, "高级数据库对象全景", "索引、视图、触发器、存储过程和递归 CTE 共同支撑性能与一致性")
    d = ImageDraw.Draw(img, "RGBA")
    panels = [
        ("19", "索引", ["idx_images_hash", "idx_images_active", "idx_tags_name_trgm", "idx_ai_desc_trgm"], BLUE),
        ("4", "视图", ["v_active_images", "v_directory_stats", "v_image_search", "v_tag_stats"], GREEN),
        ("5", "触发器", ["trg_image_after_insert", "trg_image_before_update", "trg_tag_after_delete"], CYAN),
        ("5", "存储过程", ["sp_batch_rename", "sp_restore_version", "sp_batch_insert_tags"], ORANGE),
        ("4", "函数", ["fn_log_image_insert", "fn_log_image_update", "fn_log_tag_change"], BLUE),
        ("CTE", "递归查询", ["目录树统计", "当前目录子树过滤", "目录空间报表"], GREEN),
    ]
    coords = [(78, 190), (556, 190), (1034, 190), (78, 510), (556, 510), (1034, 510)]
    for (num, name, rows, color), (x, y) in zip(panels, coords):
        shadowed_round(img, (x, y, x + 420, y + 246), 30, fill=(255, 255, 255, 224), shadow=18)
        d.rounded_rectangle((x + 26, y + 30, x + 116, y + 118), radius=24, fill=(*color, 30))
        draw_text(d, (x + 44, y + 48), num, 30 if num != "CTE" else 27, color, True, latin=True)
        draw_text(d, (x + 140, y + 36), name, 30, INK, True)
        d.rounded_rectangle((x + 140, y + 84, x + 280, y + 92), radius=4, fill=color)
        yy = y + 122
        for row in rows:
            d.ellipse((x + 42, yy + 10, x + 52, yy + 20), fill=color)
            draw_text(d, (x + 66, yy), row, 18, (40, 69, 108), True, max_width=300, line_gap=4, latin=all(ord(ch) < 128 for ch in row))
            yy += 36
    draw_text(d, (96, 800), "答辩讲法：先讲对象存在，再讲对象服务的数据库目标：查询性能、统一外模式、自动日志、批处理事务和树形统计。", 23, INK, True, max_width=1320)


def draw_directory(img: Image.Image, slide: dict[str, Any]) -> None:
    title_block(img, slide["title"], slide["subtitle"])
    paste_asset(img, "real_03_主界面_1440x900_宽屏窗口.png", (76, 190, 1040, 800), radius=30, fit="cover")
    bullet_panel(img, (1090, 204, 1490, 430), "目录树设计", ["directories.parent_id 自引用", "懒加载避免一次性扫描 UI 卡顿", "递归 CTE 统计路径与子树空间"], accent=BLUE)
    bullet_panel(img, (1090, 470, 1490, 760), "数据库收益", ["目录路径可查询", "图片与目录外键绑定", "搜索可限定当前目录及子目录"], accent=GREEN)


def draw_pipeline(img: Image.Image) -> None:
    title_block(img, "图片入库链路：缩略图、元数据与 bytea", "扫描目录后把文件系统信息结构化，缩略图缓存直接服务主界面渲染")
    steps = [
        ("磁盘目录", "遍历 JPG/PNG/GIF/BMP"),
        ("元数据抽取", "文件名、大小、格式、分辨率、hash"),
        ("缩略图缓存", "thumbnail_data bytea"),
        ("结构化入库", "images + directories"),
        ("界面复用", "网格渲染与搜索"),
    ]
    d = ImageDraw.Draw(img, "RGBA")
    y = 380
    for i, (name, desc) in enumerate(steps):
        x = 82 + i * 298
        shadowed_round(img, (x, y, x + 230, y + 170), 28, fill=(255, 255, 255, 222), shadow=18)
        d.ellipse((x + 78, y - 42, x + 152, y + 32), fill=(28, 115, 255, 225))
        draw_text(d, (x + 101, y - 26), str(i + 1), 26, (255, 255, 255), True, latin=True)
        draw_text(d, (x + 28, y + 44), name, 25, INK, True, max_width=174)
        draw_text(d, (x + 28, y + 90), desc, 18, MUTED, max_width=176)
        if i < len(steps) - 1:
            d.line((x + 236, y + 84, x + 290, y + 84), fill=(28, 115, 255, 190), width=5)
            d.polygon([(x + 290, y + 84), (x + 276, y + 74), (x + 276, y + 94)], fill=(28, 115, 255, 190))
    paste_asset(img, "real_09_图片查看器_960x680.png", (104, 604, 690, 808), radius=24, fit="cover")
    bullet_panel(img, (760, 612, 1485, 808), "为什么要把缩略图放进数据库", ["减少重复解码与磁盘读取", "搜索结果可快速回显", "图片元数据和缓存生命周期统一管理"], accent=GREEN)


def draw_version(img: Image.Image) -> None:
    title_block(img, "版本历史：从编辑行为回到数据库一致性", "image_versions 记录快照，sp_restore_version 支撑恢复，operation_logs 留痕")
    paste_asset(img, "real_14_图片编辑器_1250x938.png", (78, 190, 1080, 800), radius=30, fit="cover")
    bullet_panel(img, (1110, 210, 1490, 430), "版本表", ["image_versions 记录版本号", "保留版本文件路径与缩略图", "is_current 标记当前版本"], accent=BLUE)
    bullet_panel(img, (1110, 470, 1490, 752), "恢复链路", ["sp_restore_version 切换状态", "磁盘文件与数据库同步", "operation_logs 保存审计痕迹"], accent=GREEN)


def draw_gallery(img: Image.Image, slide: dict[str, Any]) -> None:
    title_block(img, slide["title"], slide["subtitle"])
    boxes = [(70, 180, 522, 374), (574, 180, 1026, 374), (1078, 180, 1530, 374), (70, 430, 522, 804), (574, 430, 1026, 804), (1078, 430, 1530, 804)]
    labels = ["首次启动", "数据库初始化", "图片查看", "幻灯片播放", "图片编辑器", "批量重命名"]
    for name, box, label in zip(slide["assets"], boxes, labels):
        paste_asset(img, name, box, radius=24, fit="cover")
        d = ImageDraw.Draw(img, "RGBA")
        d.rounded_rectangle((box[0] + 18, box[1] + 16, box[0] + 158, box[1] + 50), radius=14, fill=(255, 255, 255, 224))
        draw_text(d, (box[0] + 34, box[1] + 22), label, 17, INK, True)


def draw_demo(img: Image.Image) -> None:
    title_block(img, "现场演示路线", "先证明数据库初始化，再展示落库、查询、版本和幻灯片主流程")
    steps = [
        ("1", "初始化", "数据库向导检测连接并执行 schema.sql"),
        ("2", "入库", "选择目录，目录树懒加载，缩略图写入 bytea"),
        ("3", "查询", "关键词 / NL2SQL 走只读视图"),
        ("4", "恢复", "编辑图片后用版本历史回退"),
        ("5", "展示", "幻灯片播放与背景音乐收尾"),
    ]
    d = ImageDraw.Draw(img, "RGBA")
    for i, (num, name, desc) in enumerate(steps):
        x = 86 + i * 300
        y = 316 + (36 if i % 2 else 0)
        d.line((x + 112, 404, x + 300, 404), fill=(28, 115, 255, 155), width=5)
        shadowed_round(img, (x, y, x + 238, y + 198), 30, fill=(255, 255, 255, 224), shadow=20)
        d.ellipse((x + 28, y + 28, x + 84, y + 84), fill=BLUE)
        draw_text(d, (x + 47, y + 40), num, 24, (255, 255, 255), True, latin=True)
        draw_text(d, (x + 104, y + 34), name, 28, INK, True)
        draw_text(d, (x + 30, y + 108), desc, 18, MUTED, max_width=178, line_gap=5)
    paste_asset(img, "图7_演示路线.png", (210, 596, 1390, 812), radius=24, fit="contain")


def draw_summary(img: Image.Image) -> None:
    title_block(img, "总结：用数据库重新组织本地图片", "已完成主流程，云端/WebDAV、语音搜索作为后续扩展预留")
    bullet_panel(img, (86, 190, 744, 506), "已经完成并可演示", [
        "JavaFX + PostgreSQL 桌面图片管理",
        "目录树懒加载、缩略图与元数据落库",
        "AI 标签、关键词搜索、NL2SQL 只读查询",
        "版本历史、操作日志、幻灯片播放",
    ], accent=BLUE)
    bullet_panel(img, (856, 190, 1514, 506), "数据库课程价值", [
        "13 表 / 19 索引 / 4 视图",
        "5 触发器 / 5 存储过程 / 递归 CTE",
        "HikariCP、PreparedStatement、事务回滚",
        "初始化向导与离线降级保证可交付",
    ], accent=GREEN)
    bullet_panel(img, (300, 586, 1300, 792), "诚实边界", [
        "WebDAV 与云端同步只作为扩展接口，不写成完整主流程",
        "语音搜索、全网全盘检索、SQLite 多底层切换均列为未来演进",
    ], accent=ORANGE)


def render_slide(slide: dict[str, Any]) -> Image.Image:
    img = add_wash(load_bg(slide["id"]), opacity=34)
    visual = slide["visual"]
    if visual == "cover":
        draw_cover(img, slide)
    elif visual == "score":
        draw_score(img)
    elif visual == "asset_focus":
        draw_asset_focus(img, slide)
    elif visual == "tables":
        draw_tables(img)
    elif visual == "directory":
        draw_directory(img, slide)
    elif visual == "pipeline":
        draw_pipeline(img)
    elif visual == "version":
        draw_version(img)
    elif visual == "gallery":
        draw_gallery(img, slide)
    elif visual == "demo":
        draw_demo(img)
    elif visual == "summary":
        draw_summary(img)
    return img.convert("RGB")


def build_pages() -> None:
    ensure_dirs()
    for slide in SLIDES:
        out = OUTPUT_DIR / f"{slide['id']}.png"
        render_slide(slide).save(out, quality=95)
        print(f"Wrote {out}")


def build_ppt() -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]
    for slide in SLIDES:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(str(OUTPUT_DIR / f"{slide['id']}.png"), 0, 0, width=prs.slide_width, height=prs.slide_height)
    PPTX_RESCUE.parent.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_RESCUE)
    PPTX_FINAL.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(PPTX_RESCUE, PPTX_FINAL)
    print(f"Wrote {PPTX_RESCUE}")
    print(f"Copied {PPTX_FINAL}")


def build_review_html() -> None:
    if not REVIEW_SHELL.exists():
        print(f"Review shell missing: {REVIEW_SHELL}")
        return
    shell = REVIEW_SHELL.read_text(encoding="utf-8")
    slides = []
    for slide in SLIDES:
        rel = (OUTPUT_DIR / f"{slide['id']}.png").relative_to(WORK_DIR).as_posix()
        slides.append({
            "id": slide["id"],
            "code": slide["id"],
            "title": slide["title"],
            "role": slide["role"],
            "image": rel,
        })
    repl = "const sampleSlides = " + json.dumps(slides, ensure_ascii=False, indent=6) + ";"
    shell = re.sub(r"const sampleSlides = \[.*?\];", repl, shell, flags=re.S)
    REVIEW_HTML.write_text(shell, encoding="utf-8")
    print(f"Wrote {REVIEW_HTML}")


def inspect_pages() -> None:
    issues: list[str] = []
    for slide in SLIDES:
        path = OUTPUT_DIR / f"{slide['id']}.png"
        if not path.exists():
            issues.append(f"{slide['id']}: missing page image")
            continue
        with Image.open(path) as img:
            if img.size != (W, H):
                issues.append(f"{slide['id']}: size {img.size}")
            stat = ImageStatLite(img)
            if stat.stddev < 8:
                issues.append(f"{slide['id']}: image looks nearly blank")
    report = {
        "page_count": len(SLIDES),
        "size": [W, H],
        "issues": issues,
        "pptx_rescue": str(PPTX_RESCUE),
        "pptx_final": str(PPTX_FINAL),
    }
    (WORK_DIR / "rescue_deck_validation.json").write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps(report, ensure_ascii=False, indent=2))


class ImageStatLite:
    def __init__(self, img: Image.Image) -> None:
        small = img.convert("L").resize((80, 45))
        pixels = list(small.getdata())
        mean = sum(pixels) / len(pixels)
        self.stddev = math.sqrt(sum((p - mean) ** 2 for p in pixels) / len(pixels))


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--prepare", action="store_true")
    parser.add_argument("--render", action="store_true")
    parser.add_argument("--ppt", action="store_true")
    parser.add_argument("--review", action="store_true")
    parser.add_argument("--inspect", action="store_true")
    parser.add_argument("--all", action="store_true")
    args = parser.parse_args()
    if args.prepare or args.all:
        write_planning_files()
        write_prompts()
    if args.render or args.all:
        build_pages()
    if args.ppt or args.all:
        build_ppt()
    if args.review or args.all:
        build_review_html()
    if args.inspect or args.all:
        inspect_pages()


if __name__ == "__main__":
    main()
