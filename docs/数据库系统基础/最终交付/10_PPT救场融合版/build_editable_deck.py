from __future__ import annotations

import json
import shutil
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[4]
FINAL_DIR = ROOT / "docs" / "数据库系统基础" / "最终交付"
OFFICIAL = FINAL_DIR / "03_答辩PPT" / "第07组毕振岚-数据库课程设计答辩PPT.pptx"
OUT_DIR = FINAL_DIR / "10_PPT救场融合版"
EDITABLE = OUT_DIR / "第07组毕振岚-数据库课程设计答辩PPT_可编辑组件版.pptx"
VALIDATION = OUT_DIR / "editable_deck_validation.json"
FIG_DIR = FINAL_DIR / "06_证据与清单" / "figures"
SHOT_DIR = FINAL_DIR / "07_Gemini交付包" / "01_Gemini_最终报告包" / "上传资料" / "05_界面截图与图表"


W, H = Inches(13.333), Inches(7.5)

COLORS = {
    "bg": RGBColor(246, 250, 255),
    "panel": RGBColor(255, 255, 255),
    "ink": RGBColor(19, 38, 64),
    "muted": RGBColor(96, 112, 132),
    "blue": RGBColor(32, 112, 202),
    "cyan": RGBColor(37, 181, 198),
    "green": RGBColor(42, 161, 122),
    "line": RGBColor(211, 226, 244),
    "soft_blue": RGBColor(226, 241, 255),
    "soft_green": RGBColor(226, 246, 239),
    "soft_cyan": RGBColor(223, 248, 252),
    "soft_gray": RGBColor(238, 243, 249),
    "dark": RGBColor(20, 34, 55),
}


def rgb(hex_value: str) -> RGBColor:
    hex_value = hex_value.strip("#")
    return RGBColor(int(hex_value[:2], 16), int(hex_value[2:4], 16), int(hex_value[4:], 16))


def set_font(run, size=18, color=None, bold=False, font="Microsoft YaHei") -> None:
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def add_text(slide, text, x, y, w, h, size=18, color=None, bold=False, align=PP_ALIGN.LEFT, font="Microsoft YaHei"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size=size, color=color or COLORS["ink"], bold=bold, font=font)
    return box


def add_multiline(slide, lines, x, y, w, h, size=14, color=None, bullet=False, line_spacing=1.08):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = "Microsoft YaHei"
        p.line_spacing = line_spacing
        if bullet:
            p.text = "· " + line
        for run in p.runs:
            set_font(run, size=size, color=color or COLORS["muted"])
    return box


def add_rect(slide, x, y, w, h, fill=None, line=None, radius=True, transparency=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill or COLORS["panel"]
    if transparency:
        shp.fill.transparency = transparency
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp


def add_line(slide, x1, y1, x2, y2, color=None, width=1.25, arrow=False):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color or COLORS["line"]
    line.line.width = Pt(width)
    if arrow:
        line.line.end_arrowhead = True
    return line


def add_title(slide, title, subtitle=None, section="数据库系统课程设计答辩", idx=None):
    add_text(slide, section, 0.58, 0.26, 3.4, 0.25, size=8.5, color=COLORS["blue"], bold=True)
    add_text(slide, title, 0.58, 0.55, 7.6, 0.45, size=23, color=COLORS["ink"], bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.6, 1.02, 9.2, 0.28, size=11, color=COLORS["muted"])
    add_rect(slide, 0.58, 1.36, 1.1, 0.045, fill=COLORS["blue"], radius=False)
    add_rect(slide, 1.72, 1.36, 0.54, 0.045, fill=COLORS["cyan"], radius=False)
    if idx:
        add_text(slide, f"{idx:02d}", 12.05, 0.39, 0.55, 0.28, size=10, color=COLORS["muted"], bold=True, align=PP_ALIGN.RIGHT)


def add_footer(slide, idx):
    add_text(slide, "第07组 | 基于 PostgreSQL 的数字图像集成管理系统", 0.58, 7.04, 5.8, 0.18, size=7.5, color=COLORS["muted"])
    add_text(slide, f"{idx:02d} / 14", 11.98, 7.04, 0.8, 0.18, size=7.5, color=COLORS["muted"], align=PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body=None, fill=None, accent=None, title_size=13, body_size=9.5):
    add_rect(slide, x, y, w, h, fill=fill or COLORS["panel"], line=COLORS["line"])
    if accent:
        add_rect(slide, x, y, 0.08, h, fill=accent, radius=False)
    add_text(slide, title, x + 0.16, y + 0.14, w - 0.28, 0.25, size=title_size, color=COLORS["ink"], bold=True)
    if body:
        if isinstance(body, str):
            body = [body]
        add_multiline(slide, body, x + 0.16, y + 0.48, w - 0.28, h - 0.56, size=body_size, color=COLORS["muted"])


def add_metric(slide, x, y, label, value, color):
    add_rect(slide, x, y, 1.5, 0.72, fill=COLORS["panel"], line=COLORS["line"])
    add_text(slide, value, x + 0.12, y + 0.11, 0.58, 0.28, size=17, color=color, bold=True)
    add_text(slide, label, x + 0.12, y + 0.43, 1.18, 0.18, size=7.5, color=COLORS["muted"])


def add_image_fit(slide, path: Path, x, y, w, h, border=True):
    if not path.exists():
        add_card(slide, x, y, w, h, path.name, "素材缺失", fill=COLORS["soft_gray"], accent=COLORS["blue"])
        return None
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    px, py = x + (w - pw) / 2, y + (h - ph) / 2
    if border:
        add_rect(slide, x, y, w, h, fill=COLORS["panel"], line=COLORS["line"])
    pic = slide.shapes.add_picture(str(path), Inches(px), Inches(py), width=Inches(pw), height=Inches(ph))
    return pic


def add_table_shape(slide, x, y, w, h, rows, header_fill=COLORS["soft_blue"]):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = value
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if i == 0 else COLORS["panel"]
            for p in cell.text_frame.paragraphs:
                p.font.name = "Microsoft YaHei"
                p.font.size = Pt(8.2 if i else 8.6)
                p.font.bold = i == 0
                p.font.color.rgb = COLORS["ink"] if i == 0 else COLORS["muted"]
    return table_shape


def add_relation_grid(slide, x, y, w, title, rows):
    add_text(slide, title, x, y, w, 0.22, size=10.5, color=COLORS["ink"], bold=True)
    col_w = [w * 0.50, w * 0.16, w * 0.34]
    col_x = [x, x + col_w[0], x + col_w[0] + col_w[1]]
    headers = ["外键引用", "基数", "说明"]
    header_y = y + 0.30
    row_h = 0.40
    for i, header in enumerate(headers):
        add_rect(slide, col_x[i], header_y, col_w[i], 0.34, fill=COLORS["soft_blue"], line=COLORS["line"], radius=False)
        add_text(slide, header, col_x[i] + 0.03, header_y + 0.06, col_w[i] - 0.06, 0.12, size=7.4, color=COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        yy = header_y + 0.34 + r * row_h
        fill = COLORS["panel"] if r % 2 == 0 else COLORS["soft_gray"]
        for i, text in enumerate(row):
            add_rect(slide, col_x[i], yy, col_w[i], row_h, fill=fill, line=COLORS["line"], radius=False)
            size = 6.35 if i == 0 else 7.0
            add_text(slide, text, col_x[i] + 0.03, yy + 0.07, col_w[i] - 0.06, 0.14, size=size, color=COLORS["muted"], align=PP_ALIGN.CENTER if i == 1 else PP_ALIGN.LEFT)


def add_code(slide, code, x, y, w, h, size=8.8):
    add_rect(slide, x, y, w, h, fill=COLORS["dark"], line=rgb("2B4A6B"), radius=True)
    box = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.12), Inches(w - 0.24), Inches(h - 0.22))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = code
    set_font(run, size=size, color=RGBColor(229, 243, 255), font="Consolas")
    return box


def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["bg"]
    add_rect(slide, 10.7, -0.2, 2.8, 2.1, fill=COLORS["soft_cyan"], radius=True, transparency=18)
    add_rect(slide, -0.32, 5.85, 2.2, 1.6, fill=COLORS["soft_blue"], radius=True, transparency=28)


def build_deck() -> None:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    def slide(idx, title=None, subtitle=None):
        s = prs.slides.add_slide(blank)
        set_bg(s)
        if title:
            add_title(s, title, subtitle, idx=idx)
            add_footer(s, idx)
        return s

    # S01
    s = slide(1)
    add_text(s, "基于 PostgreSQL 的", 0.72, 0.78, 6.0, 0.38, size=18, color=COLORS["blue"], bold=True)
    add_text(s, "数字图像集成管理系统", 0.68, 1.18, 7.4, 0.72, size=32, color=COLORS["ink"], bold=True)
    add_text(s, "数据库系统课程设计答辩 | JavaFX + PostgreSQL 桌面端图片数据资产管理", 0.72, 2.05, 7.9, 0.28, size=12, color=COLORS["muted"])
    add_text(s, "第07组  毕振岚 | 指导教师：郭玉彬 | 2026年6月", 0.72, 2.42, 6.8, 0.25, size=10.5, color=COLORS["muted"])
    for x, y, label, value, color in [
        (0.74, 3.18, "tables", "13", COLORS["blue"]),
        (2.42, 3.18, "indexes", "19", COLORS["cyan"]),
        (4.10, 3.18, "views", "4", COLORS["green"]),
        (5.78, 3.18, "triggers/procs", "10", COLORS["blue"]),
    ]:
        add_metric(s, x, y, label, value, color)
    add_rect(s, 8.2, 0.72, 4.25, 5.55, fill=COLORS["panel"], line=COLORS["line"])
    add_rect(s, 8.45, 1.02, 3.72, 0.36, fill=COLORS["soft_blue"], line=COLORS["line"])
    add_text(s, "ImageManager 主界面", 8.62, 1.11, 2.6, 0.14, size=7.5, color=COLORS["muted"])
    add_image_fit(s, SHOT_DIR / "real_02_主界面_1200x800_默认窗口.png", 8.45, 1.5, 3.72, 2.45, border=False)
    add_card(s, 8.52, 4.25, 1.02, 0.92, "AI 标签", ["落库", "可检索"], fill=COLORS["soft_cyan"], accent=COLORS["cyan"], title_size=9, body_size=7.2)
    add_card(s, 9.74, 4.25, 1.02, 0.92, "版本", ["快照", "回滚"], fill=COLORS["soft_green"], accent=COLORS["green"], title_size=9, body_size=7.2)
    add_card(s, 10.96, 4.25, 1.02, 0.92, "CTE", ["目录树", "统计"], fill=COLORS["soft_blue"], accent=COLORS["blue"], title_size=9, body_size=7.2)
    add_footer(s, 1)

    # S02
    s = slide(2, "答辩先对齐得分点", "数据库设计是主线，功能、后台和界面都服务于落库与查询")
    score_items = [
        ("数据库设计", "30", "ER、关系模式、物理结构、索引、视图、触发器、存储过程", COLORS["blue"]),
        ("功能设计", "20", "目录、缩略图、AI 标签、NL2SQL、版本、幻灯片", COLORS["cyan"]),
        ("后台程序设计", "15", "HikariCP、PreparedStatement、事务回滚、后台扫描", COLORS["green"]),
        ("界面设计", "15", "初始化向导、主界面、编辑器、设置页、演示路径", COLORS["blue"]),
        ("报告表述", "10", "模板章节、数据字典、测试截图、总结分析", COLORS["cyan"]),
        ("PPT表达", "10", "15分钟内直入数据库含金量", COLORS["green"]),
    ]
    for i, (name, pts, body, color) in enumerate(score_items):
        x = 0.72 + (i % 2) * 6.05
        y = 1.75 + (i // 2) * 1.45
        add_card(s, x, y, 5.55, 1.08, name, body, fill=COLORS["panel"], accent=color, title_size=13, body_size=9)
        add_text(s, pts, x + 4.58, y + 0.2, 0.55, 0.25, size=18, color=color, bold=True, align=PP_ALIGN.RIGHT)
        add_text(s, "分", x + 5.12, y + 0.25, 0.25, 0.16, size=8, color=COLORS["muted"])

    # S03
    s = slide(3, "系统功能结构：图片管理不是文件夹外壳", "目录、图片、标签、AI、版本、日志共同组成可查询的数据资产")
    add_card(s, 4.85, 2.45, 3.1, 1.18, "PostgreSQL 数据资产层", ["约束、事务、索引、视图、触发器、存储过程"], fill=COLORS["soft_blue"], accent=COLORS["blue"], title_size=13, body_size=8.5)
    modules = [
        (0.8, 1.65, "目录树", "递归 CTE / 懒加载"),
        (0.8, 3.35, "图片元数据", "路径、大小、hash、bytea 缩略图"),
        (4.85, 1.05, "AI 标签", "分析结果入库后再检索"),
        (8.9, 1.65, "NL2SQL 搜索", "只读视图 + SQL 安全校验"),
        (8.9, 3.35, "版本历史", "编辑快照、恢复过程、操作日志"),
        (4.85, 4.72, "界面演示", "向导、主界面、编辑器、幻灯片"),
    ]
    for x, y, title, body in modules:
        add_card(s, x, y, 3.05, 0.95, title, body, fill=COLORS["panel"], accent=COLORS["cyan"], title_size=12, body_size=8.5)
        add_line(s, x + 1.5, y + (0.95 if y < 2.4 else 0), 6.4, 3.02, color=COLORS["line"], width=1.1)

    # S04
    s = slide(4, "ER 关系：按外键方向读，不按箭头猜", "统一口径：子表.FK -> 父表.PK；基数直接写明，避免把一对多和多对多混在一起")
    add_card(s, 0.72, 1.54, 11.78, 0.62, "读图规则", [
        "箭头或表格都按“外键引用方向”理解；真正的一对多写成 N:1；images 与 tags 的多对多必须通过 image_tags 拆成两条 N:1。"
    ], fill=COLORS["soft_blue"], accent=COLORS["blue"], title_size=11.2, body_size=8.0)

    left_rows = [
        ("directories.parent_id -> directories.id", "N:1", "目录树自引用；根目录为 NULL"),
        ("images.directory_id -> directories.id", "N:1", "一个目录可有多张图片"),
        ("operation_logs.image_id -> images.id", "0..N:1", "图片物理删除后 FK 置 NULL"),
        ("image_versions.image_id -> images.id", "N:1", "一张图片多个版本，is_current 标当前"),
        ("ai_analysis_results.image_id -> images.id", "0..1:1", "UNIQUE；一图最多一条 AI 完整结果"),
    ]
    right_rows = [
        ("tags.category_id -> tag_categories.id", "N:1", "一个分类下多个标签"),
        ("image_tags.image_id -> images.id", "N:1", "中间表引用图片"),
        ("image_tags.tag_id -> tags.id", "N:1", "中间表引用标签"),
        ("image_edit_operations.version_id -> image_versions.id", "N:1", "扩展预留：操作回放明细"),
        ("cloud_images.source_id -> cloud_sources.id", "N:1", "扩展预留：云端缓存"),
    ]
    add_relation_grid(s, 0.72, 2.42, 5.85, "图片主线与审计关系", left_rows)
    add_relation_grid(s, 6.80, 2.42, 5.75, "标签体系与扩展关系", right_rows)

    add_card(s, 0.78, 5.92, 11.55, 0.66, "答辩一句话", [
        "images 是主表；directories 是目录维度；版本、AI 结果、日志依附图片；tag_categories/tags 是标签字典；image_tags 把 images 与 tags 的 M:N 关系拆成两条 N:1。app_settings 和 search_history 没有外键。"
    ], fill=COLORS["soft_green"], accent=COLORS["green"], title_size=11.2, body_size=7.4)

    # S05
    s = slide(5, "物理结构：13 张表按职责分层", "基础对象、AI 标签、历史审计、扩展预留边界清晰")
    for x, label, value, color in [(0.72, "表", "13", COLORS["blue"]), (2.38, "索引", "19", COLORS["cyan"]), (4.04, "视图", "4", COLORS["green"]), (5.70, "触发器", "5", COLORS["blue"]), (7.36, "存储过程", "5", COLORS["cyan"])]:
        add_metric(s, x, 1.62, label, value, color)
    rows = [
        ["层次", "表/对象", "答辩时说明的职责"],
        ["目录与图片", "directories, images", "把文件系统层级、图片元数据、缩略图 bytea 结构化"],
        ["标签与 AI", "tag_categories, tags, image_tags, ai_analysis_results", "AI 识别结果先落库，再支撑多维检索"],
        ["历史与审计", "image_versions, image_edit_operations, operation_logs", "编辑快照、版本恢复、操作留痕"],
        ["运行配置", "app_settings, search_history", "保存设置与搜索轨迹"],
        ["扩展预留", "cloud_sources, cloud_images", "仅说明为 WebDAV/云端同步的未来扩展接口"],
    ]
    add_table_shape(s, 0.72, 2.65, 11.9, 3.35, rows)

    # S06
    s = slide(6, "高级数据库对象全景", "索引、视图、触发器、存储过程和递归 CTE 共同支撑性能与一致性")
    cards = [
        ("视图", ["v_active_images", "v_directory_stats", "v_image_search", "v_tag_stats"], COLORS["blue"]),
        ("索引", ["idx_images_directory_id", "idx_tags_name_trgm", "idx_ai_desc_trgm", "idx_versions_current"], COLORS["cyan"]),
        ("触发器/函数", ["trg_image_after_insert", "trg_image_before_update", "fn_log_tag_change"], COLORS["green"]),
        ("存储过程", ["sp_monthly_report", "sp_batch_rename", "sp_restore_version", "sp_batch_insert_tags"], COLORS["blue"]),
    ]
    for i, (title, lines, color) in enumerate(cards):
        add_card(s, 0.72 + i * 3.05, 1.72, 2.72, 2.05, title, lines, fill=COLORS["panel"], accent=color, title_size=13, body_size=7.8)
    add_code(s, "WITH RECURSIVE dir_tree AS (\n  SELECT directory_id, parent_id, path\n  FROM directories WHERE parent_id IS NULL\n  UNION ALL\n  SELECT d.directory_id, d.parent_id, d.path\n  FROM directories d JOIN dir_tree t\n    ON d.parent_id = t.directory_id\n)\nSELECT * FROM dir_tree;", 1.08, 4.42, 5.65, 1.82, size=8.2)
    add_card(s, 7.15, 4.42, 4.88, 1.82, "答辩抓手", ["不要只说“用了数据库”。要指出每类 SQL 对象怎样服务界面性能、搜索安全和一致性。", "老师追问时，优先打开 schema.sql 对应对象。"], fill=COLORS["soft_green"], accent=COLORS["green"], title_size=12.5, body_size=8.5)

    # S07
    s = slide(7, "完整磁盘目录树：懒加载 + 递归 CTE", "前台只展开当前需要的节点，数据库保留完整目录层级与路径统计")
    add_card(s, 0.78, 1.68, 3.2, 4.65, "目录树状态", None, fill=COLORS["panel"], accent=COLORS["cyan"])
    tree_lines = [("D:/Pictures", 0), ("2026", 1), ("课程设计截图", 2), ("相册备份", 2), ("旅行", 1), ("待整理", 1)]
    for i, (name, level) in enumerate(tree_lines):
        y = 2.15 + i * 0.48
        add_rect(s, 1.02 + level * 0.25, y, 0.13, 0.13, fill=COLORS["cyan"], radius=True)
        add_text(s, name, 1.22 + level * 0.25, y - 0.04, 2.1, 0.18, size=9, color=COLORS["ink"] if i == 0 else COLORS["muted"])
    add_code(s, "DirectoryTreeView\n  -> expand node\n  -> query children by parent_id\n  -> update UI without full scan", 4.55, 1.72, 3.35, 1.35, size=9)
    add_card(s, 8.28, 1.72, 3.8, 1.35, "数据库价值", ["目录层级在 directories 中自引用存储，路径统计由 v_directory_stats 支撑。"], fill=COLORS["soft_blue"], accent=COLORS["blue"], title_size=12, body_size=8.4)
    for i, (title, body) in enumerate([
        ("懒加载", "只查当前展开节点，避免启动时遍历整盘"),
        ("递归 CTE", "需要全局树或统计时由数据库递归展开"),
        ("后台扫描", "ScanTask 在独立线程入库，不阻塞 JavaFX"),
    ]):
        add_card(s, 4.55 + i * 2.55, 3.72, 2.28, 1.35, title, body, fill=COLORS["panel"], accent=[COLORS["blue"], COLORS["cyan"], COLORS["green"]][i], title_size=11.5, body_size=8.1)

    # S08
    s = slide(8, "图片入库链路：缩略图、元数据与 bytea", "扫描目录后把文件系统信息结构化，缩略图缓存直接服务主界面渲染")
    steps = [
        ("Files.walkFileTree", "遍历目录"),
        ("ImageService", "读取元数据"),
        ("Thumbnail", "生成 bytea"),
        ("images", "写入主表"),
        ("UI Grid", "快速渲染"),
    ]
    for i, (title, body) in enumerate(steps):
        x = 0.78 + i * 2.42
        add_card(s, x, 2.15, 1.95, 1.18, title, body, fill=COLORS["panel"], accent=[COLORS["blue"], COLORS["cyan"], COLORS["green"], COLORS["blue"], COLORS["cyan"]][i], title_size=10.5, body_size=8.0)
        if i < len(steps) - 1:
            add_line(s, x + 1.95, 2.72, x + 2.35, 2.72, color=COLORS["blue"], width=1.6)
    add_table_shape(s, 0.82, 4.28, 5.75, 1.72, [
        ["字段", "来源", "用途"],
        ["file_path", "磁盘路径", "定位与去重"],
        ["file_hash", "文件内容", "重复图片识别"],
        ["thumbnail bytea", "缩略图生成", "列表页快速显示"],
    ])
    add_image_fit(s, SHOT_DIR / "real_02_主界面_1200x800_默认窗口.png", 7.15, 4.08, 4.65, 1.95, border=True)

    # S09
    s = slide(9, "后台基座：连接池、预编译与事务回滚", "HikariCP 管连接，PreparedStatement 防注入，批量操作失败即回滚")
    layers = [
        ("Controller", "JavaFX 事件、状态栏、用户反馈", COLORS["blue"]),
        ("Service", "事务边界、批量重命名、扫描任务", COLORS["cyan"]),
        ("DAO", "PreparedStatement、结果映射", COLORS["green"]),
        ("PostgreSQL", "约束、索引、视图、触发器、过程", COLORS["blue"]),
    ]
    for i, (name, body, color) in enumerate(layers):
        add_card(s, 1.05, 1.65 + i * 1.0, 4.8, 0.72, name, body, fill=COLORS["panel"], accent=color, title_size=12, body_size=7.6)
        if i < 3:
            add_line(s, 3.45, 2.37 + i * 1.0, 3.45, 2.62 + i * 1.0, color=COLORS["line"], width=1.2)
    add_card(s, 6.55, 1.65, 2.3, 1.15, "HikariCP", ["连接复用", "降低频繁建连成本"], fill=COLORS["soft_blue"], accent=COLORS["blue"])
    add_card(s, 9.25, 1.65, 2.3, 1.15, "PreparedStatement", ["参数绑定", "防 SQL 注入"], fill=COLORS["soft_cyan"], accent=COLORS["cyan"])
    add_card(s, 6.55, 3.25, 2.3, 1.15, "Transaction", ["批量操作", "失败回滚"], fill=COLORS["soft_green"], accent=COLORS["green"])
    add_card(s, 9.25, 3.25, 2.3, 1.15, "Fallback", ["AI 不可用", "本地功能可用"], fill=COLORS["soft_blue"], accent=COLORS["blue"])
    add_code(s, "try (Connection c = dataSource.getConnection()) {\n  c.setAutoCommit(false);\n  dao.batchRename(c, items);\n  c.commit();\n} catch (Exception e) {\n  rollbackQuietly(c);\n}", 6.55, 4.88, 5.0, 1.25, size=8.3)

    # S10
    s = slide(10, "AI 标签与 NL2SQL：能问，但只能安全地问", "AI 结果先落库；自然语言查询只允许走 v_image_search 只读视图")
    pipeline = [
        ("图片", "本地文件"),
        ("AI 分析", "标签/描述"),
        ("tags + results", "结构化落库"),
        ("v_image_search", "只读外模式"),
        ("搜索结果", "缩略图网格"),
    ]
    for i, (title, body) in enumerate(pipeline):
        x = 0.72 + i * 2.38
        add_card(s, x, 1.92, 1.9, 1.0, title, body, fill=COLORS["panel"], accent=[COLORS["blue"], COLORS["cyan"], COLORS["green"], COLORS["blue"], COLORS["cyan"]][i], title_size=10.5, body_size=7.8)
        if i < 4:
            add_line(s, x + 1.9, 2.42, x + 2.27, 2.42, color=COLORS["blue"], width=1.5)
    add_card(s, 0.82, 4.03, 3.25, 1.58, "自然语言入口", ["“找上个月在北京拍的、带猫的图片”", "用户表达不直接进入数据库执行"], fill=COLORS["soft_blue"], accent=COLORS["blue"], title_size=12, body_size=8.2)
    add_card(s, 4.55, 4.03, 3.25, 1.58, "SQL 安全闸门", ["只允许 SELECT", "禁止 DDL/DML", "绑定只读视图与白名单字段"], fill=COLORS["soft_green"], accent=COLORS["green"], title_size=12, body_size=8.2)
    add_card(s, 8.28, 4.03, 3.25, 1.58, "数据库得分点", ["AI 是入口，PostgreSQL 才是持久化、检索与安全边界。"], fill=COLORS["soft_cyan"], accent=COLORS["cyan"], title_size=12, body_size=8.2)

    # S11
    s = slide(11, "版本历史：从编辑行为回到数据库一致性", "image_versions 记录快照，sp_restore_version 支撑恢复，operation_logs 留痕")
    versions = [("v1 ORIGINAL", "原图入库"), ("v2 CROP", "裁剪版本"), ("v3 ANNOTATE", "标注版本")]
    for i, (title, body) in enumerate(versions):
        x = 1.0 + i * 3.55
        add_card(s, x, 2.0, 2.6, 1.1, title, body, fill=COLORS["panel"], accent=[COLORS["blue"], COLORS["cyan"], COLORS["green"]][i], title_size=12, body_size=8.2)
        if i < 2:
            add_line(s, x + 2.6, 2.55, x + 3.36, 2.55, color=COLORS["line"], width=1.4)
    add_table_shape(s, 0.86, 4.18, 5.75, 1.55, [
        ["数据库对象", "作用"],
        ["image_versions", "保存版本元数据与当前版本标记"],
        ["image_edit_operations", "预留操作回放参数，主流程以版本快照为准"],
        ["sp_restore_version", "恢复历史版本并维护一致性"],
    ])
    add_image_fit(s, SHOT_DIR / "real_14_图片编辑器_1250x938.png", 7.0, 4.0, 4.7, 1.92, border=True)

    # S12
    s = slide(12, "界面证据：核心流程都已能演示", "启动向导、主界面、图片查看、幻灯片、编辑器、数据库初始化")
    shots = [
        ("数据库初始化向导", "real_08_数据库连接与初始化向导_760x680.png"),
        ("主界面缩略图", "real_02_主界面_1200x800_默认窗口.png"),
        ("图片编辑器", "real_14_图片编辑器_1250x938.png"),
        ("幻灯片播放", "real_12_幻灯片播放_1250x875.png"),
    ]
    for i, (label, filename) in enumerate(shots):
        x = 0.72 + (i % 2) * 5.98
        y = 1.72 + (i // 2) * 2.22
        add_text(s, label, x + 0.1, y - 0.26, 2.8, 0.2, size=10.2, color=COLORS["ink"], bold=True)
        add_image_fit(s, SHOT_DIR / filename, x, y, 5.28, 1.78, border=True)

    # S13
    s = slide(13, "现场演示路线", "先证明数据库初始化，再展示落库、查询、版本和幻灯片主流程")
    steps = [
        ("1", "启动向导", "连接 PostgreSQL 并初始化 schema.sql"),
        ("2", "导入目录", "目录树懒加载与图片元数据入库"),
        ("3", "查看 schema", "展示 13 表、索引、视图、触发器、过程"),
        ("4", "AI 标签", "逐供应商验证、失败 fallback、本地可用"),
        ("5", "NL2SQL 搜索", "只读视图与安全校验"),
        ("6", "编辑/幻灯片", "版本历史、恢复与播放体验"),
    ]
    for i, (num, title, body) in enumerate(steps):
        x = 0.86 + (i % 3) * 3.95
        y = 1.72 + (i // 3) * 1.62
        add_rect(s, x, y, 0.42, 0.42, fill=[COLORS["blue"], COLORS["cyan"], COLORS["green"]][i % 3], radius=True)
        add_text(s, num, x + 0.09, y + 0.075, 0.18, 0.12, size=9, color=RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)
        add_card(s, x + 0.55, y - 0.02, 3.05, 0.94, title, body, fill=COLORS["panel"], accent=[COLORS["blue"], COLORS["cyan"], COLORS["green"]][i % 3], title_size=11, body_size=7.7)
    add_card(s, 1.12, 5.5, 10.65, 0.78, "追问准备", ["表结构、主键外键、索引用途、触发器留痕、存储过程边界、事务回滚路径都能从 schema.sql 与 DAO/Service 代码定位。"], fill=COLORS["soft_blue"], accent=COLORS["blue"], title_size=12, body_size=8.4)

    # S14
    s = slide(14, "总结：用数据库重新组织本地图片", "已完成主流程，云端/WebDAV、语音搜索作为后续扩展预留")
    add_card(s, 0.92, 1.75, 3.45, 2.2, "已经完成", ["桌面图片管理主流程", "PostgreSQL 数据库初始化", "缩略图/元数据/AI 标签落库", "版本历史与幻灯片演示"], fill=COLORS["soft_blue"], accent=COLORS["blue"], title_size=14, body_size=8.7)
    add_card(s, 4.92, 1.75, 3.45, 2.2, "数据库含金量", ["13 张表、19 个索引", "4 个视图、5 个触发器", "5 个存储过程、递归 CTE", "HikariCP + PreparedStatement + 事务"], fill=COLORS["soft_green"], accent=COLORS["green"], title_size=14, body_size=8.7)
    add_card(s, 8.92, 1.75, 3.45, 2.2, "边界清楚", ["WebDAV/云端是预留扩展", "语音搜索不作为已实现主流程", "AI 不可用时本地功能仍可运行"], fill=COLORS["soft_cyan"], accent=COLORS["cyan"], title_size=14, body_size=8.7)
    add_text(s, "谢谢老师，请批评指正", 3.55, 5.28, 6.2, 0.55, size=26, color=COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    add_footer(s, 14)

    prs.save(str(EDITABLE))
    shutil.copy2(EDITABLE, OFFICIAL)


def validate() -> dict:
    prs = Presentation(str(OFFICIAL))
    slides = []
    for idx, slide in enumerate(prs.slides, 1):
        text_shapes = sum(1 for sh in slide.shapes if getattr(sh, "has_text_frame", False) and sh.text.strip())
        pics = sum(1 for sh in slide.shapes if sh.shape_type == 13)
        shapes = len(slide.shapes)
        full_slide_pics = 0
        for sh in slide.shapes:
            if sh.shape_type == 13:
                if sh.left <= Inches(0.05) and sh.top <= Inches(0.05) and sh.width >= W - Inches(0.1) and sh.height >= H - Inches(0.1):
                    full_slide_pics += 1
        slides.append({
            "slide": idx,
            "shapes": shapes,
            "text_shapes": text_shapes,
            "pictures": pics,
            "full_slide_pictures": full_slide_pics,
        })
    result = {
        "official": str(OFFICIAL),
        "editable_copy": str(EDITABLE),
        "slide_count": len(prs.slides),
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "slides": slides,
        "ok": len(prs.slides) == 14 and all(s["text_shapes"] >= 2 for s in slides) and all(s["full_slide_pictures"] == 0 for s in slides),
    }
    VALIDATION.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
    return result


if __name__ == "__main__":
    build_deck()
    print(json.dumps(validate(), ensure_ascii=False, indent=2))
