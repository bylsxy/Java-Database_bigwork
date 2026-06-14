from __future__ import annotations

import json
import math
import os
from pathlib import Path
import re
import shutil
import zipfile

from docx import Document
from docx.enum.section import WD_ORIENT
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Inches, Pt, RGBColor
from PIL import Image, ImageDraw, ImageFont, ImageOps
from pptx import Presentation
from pptx.util import Inches as PptInches


ROOT = Path(__file__).resolve().parents[4]
FINAL_DIR = ROOT / "docs" / "数据库系统基础" / "最终交付"
STAGE_DIR = FINAL_DIR / "01_阶段文档"
REPORT_DIR = FINAL_DIR / "02_课程报告"
PPT_DIR = FINAL_DIR / "03_答辩PPT"
SPEECH_DIR = FINAL_DIR / "04_打印讲稿"
PACKAGE_DIR = FINAL_DIR / "05_源码与运行包"
EVIDENCE_DIR = FINAL_DIR / "06_证据与清单"
FIG_DIR = EVIDENCE_DIR / "figures"
PPT_VISUAL_DIR = PPT_DIR / "page_visuals"

PROJECT_NAME = "基于 PostgreSQL 的数字图像集成管理系统"
GROUP_NAME = "第07组"
PRIMARY_NAME = "第07组毕振岚"
DATE_TEXT = "2026年6月"
MEMBERS = [
    ("毕振岚", "202425220501", "组长", "数据库设计、需求审阅、文档统筹"),
    ("陈厚华", "202425220502", "组员", "界面交互、测试与演示支持"),
    ("徐阳", "202425220527", "项目经理/报告人", "核心实现、数据库与 AI 功能、最终交付整合"),
]

SCREENSHOTS = {
    "main": ROOT / "docs" / "面向对象程序与设计" / "我们的实际写作" / "界面截图" / "real_02_主界面_1200x800_默认窗口.png",
    "welcome": ROOT / "docs" / "面向对象程序与设计" / "我们的实际写作" / "界面截图" / "real_05_首次启动向导_800x775.png",
    "db_setup": ROOT / "docs" / "面向对象程序与设计" / "我们的实际写作" / "界面截图" / "real_08_数据库连接与初始化向导_760x680.png",
    "settings": ROOT / "docs" / "面向对象程序与设计" / "我们的实际写作" / "界面截图" / "real_07_系统设置_850x900.png",
    "viewer": ROOT / "docs" / "面向对象程序与设计" / "我们的实际写作" / "界面截图" / "real_10_图片查看器_1200x850.png",
    "editor": ROOT / "docs" / "面向对象程序与设计" / "我们的实际写作" / "界面截图" / "real_14_图片编辑器_1250x938.png",
    "slideshow": ROOT / "docs" / "面向对象程序与设计" / "我们的实际写作" / "界面截图" / "real_12_幻灯片播放_1250x875.png",
    "rename": ROOT / "docs" / "面向对象程序与设计" / "我们的实际写作" / "界面截图" / "real_16_批量重命名_563x450.png",
}

CORE_TABLES = [
    ("directories", "目录表", "id 主键；parent_id 自引用；dir_path 唯一", "把磁盘目录映射为树形结构，支撑递归统计和当前目录过滤"),
    ("images", "图片主表", "id 主键；directory_id 外键；thumbnail bytea；file_hash；逻辑删除", "保存图片元数据、缩略图、AI 状态，是系统 CRUD 的中心"),
    ("operation_logs", "操作日志", "image_id 外键；ON DELETE SET NULL", "记录新增、重命名、删除、标签变更和版本恢复等审计信息"),
    ("tag_categories", "标签分类", "name 唯一", "定义场景、物体、人物、名人、颜色、动作、文字等标签大类"),
    ("tags", "标签值", "category_id 外键；UNIQUE(category_id, name)", "保存可复用标签词条，避免重复字符串散落在图片记录中"),
    ("image_tags", "图片标签关联", "image_id/tag_id 双外键；UNIQUE(image_id, tag_id)", "实现图片与标签的多对多关系，保留来源和置信度字段"),
    ("ai_analysis_results", "AI 分析结果", "image_id 唯一外键；raw_response；description", "缓存模型原始 JSON、自然语言描述、人数和模型名"),
    ("search_history", "搜索历史", "query_text；search_mode；generated_sql", "记录关键词和 AI SQL 搜索，便于回溯查询行为"),
    ("image_versions", "版本历史", "image_id 外键；UNIQUE(image_id, version_num)", "保存编辑版本快照、缩略图和当前版本标记"),
    ("image_edit_operations", "编辑操作参数", "version_id 外键；parameters JSONB", "预留更细粒度的编辑参数记录"),
    ("app_settings", "应用设置", "key 主键", "保存非敏感配置，如扫描目录、欢迎页、幻灯片设置"),
    ("cloud_sources", "云端配置预留", "config JSONB", "只作为 WebDAV/网盘扩展表结构预留"),
    ("cloud_images", "云端图片缓存预留", "source_id 外键；UNIQUE(source_id, remote_path)", "只作为云端图片元数据缓存预留"),
]

REQUIREMENTS = [
    ("FR-01", "启动引导与数据库自检", "欢迎向导选择扫描目录，数据库不可用时保留本地浏览，并可进入初始化向导。"),
    ("FR-02", "目录树与缩略图浏览", "显示完整磁盘根节点，子目录懒加载；目录图片入库并展示缩略图和元数据。"),
    ("FR-03", "图片增删改查", "支持搜索、复制粘贴、删除、单张/批量重命名，数据库与磁盘文件保持一致。"),
    ("FR-04", "AI 标签与标签管理", "调用 OpenAI-compatible 视觉接口，解析 JSON，写入标签、AI 结果和操作日志。"),
    ("FR-05", "关键词与 NL2SQL 搜索", "关键词搜索覆盖文件名、目录、格式、分辨率、标签和 AI 描述；AI 搜索生成只读 SELECT。"),
    ("FR-06", "图片编辑与版本历史", "支持裁切、画笔、文字、箭头、矩形标注，每次保存产生版本快照，可恢复。"),
    ("FR-07", "幻灯片播放与音乐", "支持多选或整目录播放、自动播放、全屏、缩略图跳转和内置/自定义音乐。"),
    ("FR-08", "设置与运行配置", "配置数据库、扫描目录、AI fallback、请求间隔、单批上限、幻灯片参数。"),
]

DATA_DICTIONARY = [
    ("D01", "目录节点", "磁盘名称、文件夹名称、展开状态、完整路径", "目录树区域", "用户选择扫描范围和当前工作目录"),
    ("D02", "图片缩略项", "缩略图、文件名、格式、分辨率、大小、修改时间、AI 标记", "主界面缩略图网格", "用于快速浏览和选择图片"),
    ("D03", "图片元数据", "文件路径、目录、大小、宽高、格式、缩略图、哈希、删除状态", "数据库与图片信息面板", "数据库中 `images` 表的核心数据"),
    ("D04", "标签数据", "分类、标签名、图片关联、来源、置信度字段", "标签管理与搜索", "AI 或手动标签均参与检索"),
    ("D05", "AI 分析结果", "原始 JSON、描述、人数、模型名、分析时间", "AI 扫描与搜索", "保留模型输出，避免重复请求"),
    ("D06", "搜索记录", "查询文本、模式、生成 SQL、结果数量、搜索时间", "搜索栏与后台审计", "记录关键词和 AI SQL 搜索"),
    ("D07", "版本快照", "版本号、版本文件、尺寸、缩略图、编辑类型、当前标记", "图片编辑器版本时间轴", "支持恢复历史版本"),
    ("D08", "数据库连接配置", "JDBC URL、用户名、密码、连接检测状态", "数据库初始化向导", "只写入本机配置，不进入源码包"),
    ("D09", "幻灯片设置", "播放间隔、播放顺序、音乐、音量", "设置页与幻灯片窗口", "决定自动播放体验"),
    ("D10", "操作日志", "图片、操作类型、旧值、新值、操作时间", "数据库日志表", "由触发器与应用层共同维护"),
]

SQL_OBJECTS = {
    "视图": [
        ("v_active_images", "过滤逻辑删除后的图片，作为常用外模式"),
        ("v_directory_stats", "汇总目录图片数量和容量"),
        ("v_image_search", "聚合图片、目录、标签和 AI 描述，支撑关键词与 NL2SQL"),
        ("v_tag_stats", "统计标签使用情况"),
    ],
    "触发器": [
        ("trg_image_after_insert", "图片新增后写操作日志"),
        ("trg_image_before_update", "图片名称或删除状态变化前记录差异"),
        ("trg_image_after_delete", "图片记录删除后保留日志"),
        ("trg_tag_after_insert", "标签关联新增后写日志"),
        ("trg_tag_after_delete", "标签关联删除后写日志"),
    ],
    "存储过程": [
        ("sp_monthly_report", "按月份统计图片新增与容量变化"),
        ("sp_directory_report", "使用递归 CTE 统计目录子树"),
        ("sp_batch_rename", "数据库层批量重命名过程"),
        ("sp_restore_version", "恢复指定图片历史版本的数据库状态"),
        ("sp_batch_insert_tags", "为图片批量插入 AI 标签并自动创建标签值"),
    ],
    "索引": [
        ("idx_images_active", "部分索引，加速目录内未删除图片查询"),
        ("idx_images_hash", "唯一部分索引，加速图片哈希去重和定位"),
        ("idx_images_ai_pending", "部分索引，加速待 AI 处理图片定位"),
        ("idx_tags_name_trgm", "pg_trgm GIN，加速标签模糊匹配"),
        ("idx_ai_desc_trgm", "pg_trgm GIN，加速 AI 描述模糊匹配"),
        ("idx_versions_current", "部分索引，加速当前版本查询"),
    ],
}


def ensure_dirs() -> None:
    for path in [STAGE_DIR, REPORT_DIR, PPT_DIR, SPEECH_DIR, PACKAGE_DIR, EVIDENCE_DIR, FIG_DIR, PPT_VISUAL_DIR]:
        path.mkdir(parents=True, exist_ok=True)


def font_path(*names: str) -> str:
    candidates = []
    for name in names:
        candidates.append(Path("C:/Windows/Fonts") / name)
    candidates += [
        Path("C:/Windows/Fonts/msyh.ttc"),
        Path("C:/Windows/Fonts/simhei.ttf"),
        Path("C:/Windows/Fonts/simsun.ttc"),
    ]
    for candidate in candidates:
        if candidate.exists():
            return str(candidate)
    return ""


FONT_REGULAR = font_path("msyh.ttc")
FONT_BOLD = font_path("msyhbd.ttc", "simhei.ttf")
FONT_MONO = font_path("consola.ttf", "cascadiamono.ttf", "msyh.ttc")


def pil_font(size: int, bold: bool = False, mono: bool = False) -> ImageFont.FreeTypeFont:
    path = FONT_MONO if mono else (FONT_BOLD if bold else FONT_REGULAR)
    if path:
        return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def text_width(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.ImageFont) -> int:
    if not text:
        return 0
    box = draw.textbbox((0, 0), text, font=font)
    return box[2] - box[0]


def wrap_text(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.ImageFont, width: int) -> list[str]:
    lines: list[str] = []
    for raw in str(text).split("\n"):
        line = ""
        for ch in raw:
            test = line + ch
            if line and text_width(draw, test, font) > width:
                lines.append(line)
                line = ch
            else:
                line = test
        lines.append(line)
    return lines


def draw_wrapped(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int],
    text: str,
    font: ImageFont.ImageFont,
    fill: tuple[int, int, int],
    width: int,
    line_gap: int = 8,
) -> int:
    x, y = xy
    line_height = font.size + line_gap
    for line in wrap_text(draw, text, font, width):
        draw.text((x, y), line, font=font, fill=fill)
        y += line_height
    return y


def add_doc_font(run, size: int | None = None, bold: bool | None = None, color: str | None = None) -> None:
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.bold = bold
    run.font.name = "微软雅黑"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
    if color:
        run.font.color.rgb = RGBColor.from_string(color)


def set_cell_text(cell, text: str, bold: bool = False) -> None:
    cell.text = ""
    p = cell.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = p.add_run(str(text))
    add_doc_font(run, 9, bold)
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER


def shade_cell(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def set_table_borders(table) -> None:
    tbl = table._tbl
    tbl_pr = tbl.tblPr
    borders = OxmlElement("w:tblBorders")
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        tag = OxmlElement(f"w:{edge}")
        tag.set(qn("w:val"), "single")
        tag.set(qn("w:sz"), "6")
        tag.set(qn("w:space"), "0")
        tag.set(qn("w:color"), "B7C0C7")
        borders.append(tag)
    tbl_pr.append(borders)


def style_document(doc: Document, title: str) -> None:
    section = doc.sections[0]
    section.top_margin = Cm(2.0)
    section.bottom_margin = Cm(2.0)
    section.left_margin = Cm(2.2)
    section.right_margin = Cm(2.2)
    styles = doc.styles
    styles["Normal"].font.name = "微软雅黑"
    styles["Normal"]._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
    styles["Normal"].font.size = Pt(10.5)
    for style_name, size in [("Title", 22), ("Heading 1", 16), ("Heading 2", 13), ("Heading 3", 11)]:
        style = styles[style_name]
        style.font.name = "微软雅黑"
        style._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
        style.font.size = Pt(size)
        style.font.bold = True
        style.font.color.rgb = RGBColor(32, 58, 83)
    title_p = doc.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title_p.add_run(title)
    add_doc_font(run, 22, True, "17324D")
    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = sub.add_run(f"{GROUP_NAME}  |  {DATE_TEXT}")
    add_doc_font(r, 11, False, "4B5B66")


def add_paragraph(doc: Document, text: str, style: str | None = None) -> None:
    p = doc.add_paragraph(style=style)
    p.paragraph_format.first_line_indent = Pt(21) if style is None else None
    p.paragraph_format.line_spacing = 1.25
    run = p.add_run(text)
    add_doc_font(run, 10.5)


def add_heading(doc: Document, text: str, level: int = 1) -> None:
    doc.add_heading(text, level=level)


def add_table(doc: Document, headers: list[str], rows: list[tuple | list], widths: list[float] | None = None) -> None:
    table = doc.add_table(rows=1, cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    set_table_borders(table)
    for i, header in enumerate(headers):
        set_cell_text(table.rows[0].cells[i], header, True)
        shade_cell(table.rows[0].cells[i], "D9EAF7")
    for row in rows:
        cells = table.add_row().cells
        for i, value in enumerate(row):
            set_cell_text(cells[i], str(value))
    if widths:
        for row in table.rows:
            for idx, width in enumerate(widths):
                row.cells[idx].width = Cm(width)
    doc.add_paragraph()


def add_picture(doc: Document, image_path: Path, caption: str, width_cm: float = 14.5) -> None:
    if not image_path.exists():
        add_paragraph(doc, f"图未找到：{image_path}")
        return
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run().add_picture(str(image_path), width=Cm(width_cm))
    cap = doc.add_paragraph()
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = cap.add_run(caption)
    add_doc_font(r, 9, False, "5A6670")


def rounded_rectangle(draw: ImageDraw.ImageDraw, xy, radius: int, fill, outline=None, width: int = 1):
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def card(draw, xy, title, body, accent=(30, 128, 128), width=2):
    x1, y1, x2, y2 = xy
    rounded_rectangle(draw, xy, 22, (255, 255, 255), (188, 204, 213), 2)
    draw.rectangle((x1, y1, x1 + 12, y2), fill=accent)
    draw_wrapped(draw, (x1 + 34, y1 + 22), title, pil_font(28, True), (28, 50, 70), x2 - x1 - 54, 6)
    draw_wrapped(draw, (x1 + 34, y1 + 66), body, pil_font(20), (72, 82, 88), x2 - x1 - 54, 7)


def canvas(width=1600, height=900):
    img = Image.new("RGB", (width, height), (246, 249, 247))
    draw = ImageDraw.Draw(img)
    for i in range(0, width, 48):
        draw.line((i, 0, i, height), fill=(235, 242, 240), width=1)
    for j in range(0, height, 48):
        draw.line((0, j, width, j), fill=(235, 242, 240), width=1)
    draw.ellipse((width - 420, -180, width + 100, 340), fill=(224, 241, 238))
    draw.ellipse((-220, height - 220, 380, height + 260), fill=(249, 233, 213))
    return img, draw


def image_fit(path: Path, size: tuple[int, int]) -> Image.Image:
    if not path.exists():
        ph = Image.new("RGB", size, (226, 235, 239))
        d = ImageDraw.Draw(ph)
        d.text((30, size[1] // 2 - 15), "image missing", font=pil_font(28), fill=(70, 80, 90))
        return ph
    im = Image.open(path).convert("RGB")
    return ImageOps.contain(im, size, method=Image.Resampling.LANCZOS)


def paste_cover(im: Image.Image, src: Path, box: tuple[int, int, int, int], border=(174, 190, 198)) -> None:
    x1, y1, x2, y2 = box
    draw = ImageDraw.Draw(im)
    rounded_rectangle(draw, box, 24, (255, 255, 255), border, 2)
    inner = (x2 - x1 - 26, y2 - y1 - 26)
    pic = image_fit(src, inner)
    px = x1 + 13 + (inner[0] - pic.width) // 2
    py = y1 + 13 + (inner[1] - pic.height) // 2
    im.paste(pic, (px, py))


def save_fig(name: str, img: Image.Image) -> Path:
    path = FIG_DIR / name
    img.save(path, quality=95)
    return path


def generate_figures() -> dict[str, Path]:
    figures: dict[str, Path] = {}

    img, draw = canvas()
    draw_wrapped(draw, (90, 70), "系统功能结构图", pil_font(46, True), (24, 50, 68), 900)
    modules = [
        ("启动与数据库向导", "欢迎页、连接检测、建库、schema 初始化、离线降级"),
        ("目录与缩略图管理", "完整磁盘根、子目录懒加载、元数据入库、bytea 缩略图"),
        ("图片操作", "复制粘贴、删除、单张/批量重命名、事务回滚"),
        ("AI 标签与搜索", "AI JSON 入库、标签多对多、关键词、NL2SQL"),
        ("编辑与版本历史", "编辑器、版本快照、时间轴、恢复"),
        ("幻灯片与设置", "多选播放、音乐、主题、AI fallback 和运行参数"),
    ]
    positions = [(90, 180), (570, 180), (1050, 180), (90, 500), (570, 500), (1050, 500)]
    for (title, body), (x, y) in zip(modules, positions):
        card(draw, (x, y, x + 410, y + 230), title, body)
    figures["function"] = save_fig("图1_系统功能结构图.png", img)

    img, draw = canvas()
    draw_wrapped(draw, (90, 60), "数据库 ER 与关系模式核心", pil_font(44, True), (24, 50, 68), 1000)
    boxes = {
        "directories": (90, 190, 390, 320),
        "images": (520, 190, 850, 340),
        "operation_logs": (980, 190, 1320, 320),
        "tag_categories": (90, 520, 390, 650),
        "tags": (520, 500, 850, 650),
        "image_tags": (980, 500, 1320, 650),
        "ai_analysis_results": (980, 690, 1320, 800),
        "image_versions": (520, 690, 850, 800),
    }
    table_bodies = {
        "directories": "id PK\nparent_id -> directories\npath UNIQUE",
        "images": "id PK\ndirectory_id -> directories\nthumbnail bytea\nis_deleted",
        "operation_logs": "image_id -> images\nold/new value\noperated_at",
        "tag_categories": "id PK\nname UNIQUE\ndisplay_name",
        "tags": "category_id -> tag_categories\nUNIQUE(category, name)",
        "image_tags": "image_id -> images\ntag_id -> tags\nUNIQUE(image, tag)",
        "ai_analysis_results": "image_id UNIQUE\nraw_response\ndescription",
        "image_versions": "image_id -> images\nversion_num\nis_current",
    }
    for name, xy in boxes.items():
        card(draw, xy, name, table_bodies[name], accent=(85, 142, 121))
    def arrow(a, b):
        ax1, ay1, ax2, ay2 = boxes[a]
        bx1, by1, bx2, by2 = boxes[b]
        start = (ax2, (ay1 + ay2) // 2)
        end = (bx1, (by1 + by2) // 2)
        draw.line((start, end), fill=(70, 100, 115), width=4)
        draw.polygon([(end[0], end[1]), (end[0] - 14, end[1] - 8), (end[0] - 14, end[1] + 8)], fill=(70, 100, 115))
    arrow("directories", "images")
    arrow("images", "operation_logs")
    arrow("tag_categories", "tags")
    arrow("tags", "image_tags")
    draw.line((1150, 500, 720, 340), fill=(70, 100, 115), width=4)
    draw.polygon([(720, 340), (736, 338), (727, 353)], fill=(70, 100, 115))
    draw.line((1150, 690, 850, 340), fill=(70, 100, 115), width=4)
    draw.line((690, 690, 690, 340), fill=(70, 100, 115), width=4)
    figures["er"] = save_fig("图2_ER关系模式图.png", img)

    img, draw = canvas()
    draw_wrapped(draw, (90, 60), "总体流程：从启动到数据库化检索", pil_font(44, True), (24, 50, 68), 1000)
    steps = [
        ("启动自检", "读取外部配置\nHikariCP 连接 PostgreSQL"),
        ("引导建库", "连接失败时保留本地浏览\n向导一键创建 image_manager"),
        ("扫描入库", "目录与图片元数据\n缩略图 bytea、哈希、AI 状态"),
        ("AI 标签", "视觉模型返回 JSON\n标签与描述结构化入库"),
        ("查询/编辑", "关键词或 NL2SQL 搜索\n编辑版本与事务日志"),
    ]
    x = 80
    for i, (title, body) in enumerate(steps):
        card(draw, (x, 300, x + 260, 520), title, body, accent=(32, 128, 142))
        if i < len(steps) - 1:
            draw.line((x + 270, 410, x + 330, 410), fill=(186, 91, 40), width=5)
            draw.polygon([(x + 330, 410), (x + 312, 400), (x + 312, 420)], fill=(186, 91, 40))
        x += 310
    figures["flow"] = save_fig("图3_系统总体流程图.png", img)

    img, draw = canvas()
    draw_wrapped(draw, (90, 60), "SQL 编程对象与性能设计", pil_font(44, True), (24, 50, 68), 1000)
    object_cards = [
        ("13 张表", "核心：directories, images, image_tags, ai_analysis_results, image_versions"),
        ("19 个索引", "18 个普通索引 + 1 个唯一索引，覆盖目录过滤、哈希去重、待 AI 和模糊搜索"),
        ("4 个视图", "v_active_images, v_directory_stats, v_image_search, v_tag_stats"),
        ("5 个触发器", "图片和标签变更自动写 operation_logs"),
        ("5 个存储过程", "月报、目录递归统计、批量改名、版本恢复、批量标签"),
        ("递归 CTE", "目录子树统计和当前目录下多层图片过滤"),
    ]
    for i, (title, body) in enumerate(object_cards):
        x = 100 + (i % 3) * 480
        y = 190 + (i // 3) * 300
        card(draw, (x, y, x + 400, y + 220), title, body, accent=(186, 91, 40))
    figures["sql_objects"] = save_fig("图4_SQL对象与性能设计.png", img)

    img, draw = canvas()
    draw_wrapped(draw, (90, 60), "NL2SQL 安全查询链路", pil_font(44, True), (24, 50, 68), 1000)
    steps = [
        ("自然语言", "例如：找有海边、多人、红衣服的照片"),
        ("AI 生成 SQL", "附带 schema 和视图说明\n要求只返回 SELECT"),
        ("安全校验", "必须 SELECT\n拒绝 DROP/DELETE/UPDATE/INSERT/ALTER\n只读连接，5 秒超时"),
        ("v_image_search", "聚合图片、目录、标签、AI 描述\n结果最多 1000 行"),
        ("缩略图结果", "回查 ImageDao\n展示缩略图与元数据\n写 search_history"),
    ]
    for i, (title, body) in enumerate(steps):
        y = 170 + i * 130
        card(draw, (160, y, 560, y + 100), title, body, accent=(85, 142, 121))
        if i < len(steps) - 1:
            draw.line((360, y + 103, 360, y + 126), fill=(186, 91, 40), width=4)
            draw.polygon([(360, y + 126), (350, y + 108), (370, y + 108)], fill=(186, 91, 40))
    paste_cover(img, SCREENSHOTS["main"], (760, 190, 1430, 670))
    figures["nl2sql"] = save_fig("图5_NL2SQL安全链路.png", img)

    img, draw = canvas()
    draw_wrapped(draw, (90, 60), "后台程序设计：分层、事务与降级", pil_font(44, True), (24, 50, 68), 1100)
    layers = [
        ("FXML / Controller", "MainController、SettingsController、DatabaseSetupDialog"),
        ("Service", "SearchService、ImageServiceImpl、EditService、ScanTask、DatabaseBootstrapService"),
        ("DAO", "DirectoryDao、ImageDao、TagDao、VersionDao、SettingsDao"),
        ("PostgreSQL", "schema.sql、视图、索引、触发器、存储过程"),
    ]
    for i, (title, body) in enumerate(layers):
        y = 165 + i * 155
        card(draw, (120, y, 780, y + 115), title, body, accent=(32, 128, 142))
        if i < len(layers) - 1:
            draw.line((450, y + 118, 450, y + 150), fill=(70, 100, 115), width=4)
    notes = [
        "HikariCP 单例连接池复用连接",
        "PreparedStatement 绑定参数，避免 SQL 拼接注入",
        "批量重命名使用事务，失败时数据库和文件名一起回滚",
        "AI fallback 逐个供应商验证，配置缺失时跳过 AI 阶段",
        "数据库不可用时进入本地浏览降级，不阻塞主界面",
    ]
    for i, note in enumerate(notes):
        card(draw, (900, 170 + i * 120, 1430, 260 + i * 120), f"设计点 {i+1}", note, accent=(186, 91, 40))
    figures["backend"] = save_fig("图6_后台架构图.png", img)

    img, draw = canvas()
    draw_wrapped(draw, (90, 60), "答辩演示路线", pil_font(44, True), (24, 50, 68), 1000)
    demo = [
        ("1. 启动", "双击 exe 或运行 JAR"),
        ("2. 初始化", "数据库向导检测连接，一键建库并执行 schema.sql"),
        ("3. 入库", "选择目录，懒加载目录树，扫描图片元数据和缩略图"),
        ("4. 搜索", "关键词搜索与 AI 智能搜索，展示只读 SQL 安全链路"),
        ("5. 版本", "编辑图片，保存版本，查看时间轴并恢复"),
        ("6. 交付", "展示 JAR、portable zip、源码包和验证日志"),
    ]
    for i, (title, body) in enumerate(demo):
        x = 100 + (i % 2) * 700
        y = 170 + (i // 2) * 220
        card(draw, (x, y, x + 560, y + 150), title, body, accent=(85, 142, 121))
    figures["demo"] = save_fig("图7_演示路线.png", img)

    return figures


def create_stage_requirement(figures: dict[str, Path]) -> Path:
    doc = Document()
    style_document(doc, f"{PROJECT_NAME}系统需求分析说明书")
    add_heading(doc, "文档信息", 1)
    add_table(doc, ["项目", "内容"], [
        ("项目名称", PROJECT_NAME),
        ("文档版本", "v4.0 最终交付版"),
        ("编写与修订", "徐阳、毕振岚、陈厚华"),
        ("技术口径", "Java 21、JavaFX 21.0.6、PostgreSQL 18.3、Maven 3.9.x"),
        ("修订说明", "以当前 README、schema.sql 和源码为准，剔除旧版 JavaFX 26、SQLite 已实现、完整云端主流程等不准确表述。"),
    ], [4, 11])
    add_heading(doc, "1 引言", 1)
    add_paragraph(doc, "本文档描述基于 PostgreSQL 的数字图像集成管理系统的最终需求。系统面向个人电脑机主，目标是把分散在磁盘目录中的图片转化为可检索、可审计、可维护的数据库化资源，同时保留桌面软件的直接操作体验。")
    add_paragraph(doc, "本次最终交付将数据库作为核心能力：目录、图片元数据、缩略图、AI 标签、搜索历史、版本历史和操作日志均进入 PostgreSQL；AI、WebDAV 等网络能力按可选增强处理，其中 WebDAV 仅保留依赖和表结构扩展，不写成已完成主流程。")
    add_heading(doc, "2 系统范围与用户", 1)
    add_paragraph(doc, "系统运行在 Windows 桌面环境，使用 JavaFX 提供图形界面，使用 PostgreSQL 存储结构化数据和缩略图缓存。用户不需要了解数据库结构，只需要选择图片目录、浏览图片、执行搜索、管理标签、编辑图片和播放幻灯片。")
    add_table(doc, ["角色", "主要任务"], [
        ("机主", "选择扫描目录，浏览、搜索、删除、复制、粘贴、重命名和编辑图片。"),
        ("机主", "配置数据库连接、AI fallback、扫描参数和幻灯片播放参数。"),
        ("系统", "维护目录树、图片元数据、标签、缩略图、版本和日志的一致性。"),
        ("外部 AI 服务", "在用户配置后提供图像识别和自然语言转 SQL 能力；未配置时系统跳过 AI 阶段。"),
    ], [4, 11])
    add_heading(doc, "3 功能需求", 1)
    add_picture(doc, figures["function"], "图 1 系统功能结构图")
    add_table(doc, ["编号", "功能", "需求说明"], REQUIREMENTS, [2, 4, 9])
    add_heading(doc, "4 总体流程", 1)
    add_picture(doc, figures["flow"], "图 2 系统总体流程")
    add_paragraph(doc, "系统启动后先执行数据库连接自检。如果连接失败，用户仍可进入本地图片浏览界面，同时可以通过数据库连接与初始化向导填写连接信息、检测连接、创建 image_manager 数据库并执行 JAR 内嵌的 schema.sql。目录扫描阶段把目录、图片元数据、缩略图和 AI 处理状态写入数据库，后续搜索、标签和版本历史都围绕这些数据展开。")
    add_heading(doc, "5 用户视角数据字典", 1)
    add_table(doc, ["编号", "数据", "数据组成", "所属界面/环节", "说明"], DATA_DICTIONARY, [1.6, 2.8, 5.2, 3.3, 4.0])
    add_heading(doc, "6 非功能需求", 1)
    add_table(doc, ["类别", "要求"], [
        ("性能", "目录树必须懒加载，缩略图生成和数据库检查尽量放入后台任务；AI 扫描不得阻塞 UI。"),
        ("安全", "删除有二次确认；NL2SQL 只允许 SELECT，拒绝危险关键字；AI 密钥和数据库密码不进入源码包。"),
        ("可靠性", "批量重命名使用事务，失败时回滚数据库和已改动文件名；数据库不可用时保留本地浏览降级。"),
        ("可维护性", "Controller/Service/DAO 分层，数据库对象集中在 schema.sql，配置文件和资源文件分离。"),
        ("扩展性", "标签分类、AI endpoint、视图和云端预留表支持后续扩展，但本次验收只承诺本地图片管理主流程。"),
    ], [3, 12])
    add_heading(doc, "7 需求跟踪", 1)
    add_table(doc, ["需求", "设计/实现证据", "覆盖状态"], [
        ("目录和图片管理", "MainController、ImageServiceImpl、DirectoryDaoImpl、ImageDaoImpl、images/directories 表", "已覆盖"),
        ("AI 标签与搜索", "OpenAICompatibleService、ScanTask、TagDaoImpl、image_tags、ai_analysis_results、v_image_search", "已覆盖，需外部 AI 配置"),
        ("版本历史", "EditService、VersionDaoImpl、image_versions 表、sp_restore_version", "已覆盖"),
        ("数据库初始化", "DatabaseSetupDialog、DatabaseBootstrapService、sql/schema.sql", "已覆盖"),
        ("云端/WebDAV", "pom.xml Sardine 依赖、cloud_sources/cloud_images 表", "仅预留扩展，不作为本次主流程"),
    ], [4, 8, 3])
    out = STAGE_DIR / f"{PRIMARY_NAME}-系统需求分析说明书.docx"
    doc.save(out)
    return out


def create_stage_overview(figures: dict[str, Path]) -> Path:
    doc = Document()
    style_document(doc, f"{PROJECT_NAME}概要设计说明书")
    add_heading(doc, "1 设计目标", 1)
    add_paragraph(doc, "概要设计的目标是把需求阶段提出的图片管理、AI 标签、智能搜索、版本历史和幻灯片播放转化为可实现的系统结构和数据库结构。本系统采用桌面端 JavaFX 应用形态，数据库是课程设计的核心支撑，而不是简单的附属配置表。")
    add_heading(doc, "2 总体结构", 1)
    add_picture(doc, figures["backend"], "图 1 系统分层结构")
    add_table(doc, ["层次", "主要职责", "代表文件"], [
        ("界面层", "FXML 布局、控制器事件、用户提示、窗口跳转", "MainView.fxml、MainController、SettingsController"),
        ("服务层", "目录扫描、图片操作、AI 搜索、版本保存、数据库初始化、音乐播放", "SearchService、ImageServiceImpl、EditService、ScanTask、DatabaseBootstrapService"),
        ("DAO 层", "PreparedStatement 查询和更新、事务边界、对象映射", "ImageDaoImpl、TagDaoImpl、VersionDaoImpl、DirectoryDaoImpl"),
        ("数据库层", "表、约束、索引、视图、触发器、存储过程和递归 CTE", "sql/schema.sql、sql/data.sql"),
    ], [3, 6, 6])
    add_heading(doc, "3 数据库概要设计", 1)
    add_picture(doc, figures["er"], "图 2 ER 与关系模式核心")
    add_table(doc, ["表名", "中文含义", "关键约束", "设计作用"], CORE_TABLES, [3, 3, 5, 5])
    add_heading(doc, "4 SQL 对象设计", 1)
    add_picture(doc, figures["sql_objects"], "图 3 SQL 编程对象与性能设计")
    for category, rows in SQL_OBJECTS.items():
        add_heading(doc, category, 2)
        add_table(doc, ["对象", "作用"], rows, [5, 10])
    add_heading(doc, "5 模块划分", 1)
    add_table(doc, ["模块", "设计说明"], [
        ("启动与降级模块", "App 启动时检查数据库状态，失败时打开本地浏览，并把数据库向导作为恢复入口。"),
        ("目录扫描模块", "DirectoryScanner 和 ScanTask 负责扫描文件系统、估算进度、写入目录和图片记录。"),
        ("图片管理模块", "ImageServiceImpl 处理删除、复制、粘贴、单张/批量重命名和事务回滚。"),
        ("AI 标签模块", "OpenAICompatibleService 调用兼容接口，TagDaoImpl 负责标签和 AI 结果入库。"),
        ("搜索模块", "SearchService 提供关键词和 AI SQL 两种模式，TagDaoImpl 负责 SQL 安全校验和查询。"),
        ("版本模块", "EditService 和 VersionDaoImpl 保存版本快照，数据库保留版本元数据。"),
        ("打包模块", "Maven shade 生成依赖完整 JAR，package-stable.ps1 生成 Windows portable zip/exe。"),
    ], [4, 11])
    add_heading(doc, "6 关键取舍", 1)
    add_paragraph(doc, "数据库表中保留 cloud_sources 和 cloud_images，是为了展示云端扩展设计的关系和约束，但当前前台没有完整 WebDAV 主流程。因此最终报告把它写为预留扩展，不把它列为验收主功能。")
    add_paragraph(doc, "批量重命名、版本恢复等主流程在 Java 侧使用 DAO/Service 事务完成；schema.sql 中的存储过程体现数据库层可编程能力和课程设计完整性，不夸大为所有前台操作都直接调用存储过程。")
    out = STAGE_DIR / f"{PRIMARY_NAME}-概要设计说明书.docx"
    doc.save(out)
    return out


def create_stage_detail(figures: dict[str, Path]) -> Path:
    doc = Document()
    style_document(doc, f"{PROJECT_NAME}详细设计说明书")
    add_heading(doc, "1 开发环境与技术", 1)
    add_table(doc, ["项目", "真实口径"], [
        ("语言与运行环境", "Java 21 编译目标，JDK 21 及以上运行"),
        ("界面框架", "JavaFX 21.0.6、FXML、CSS、JavaFX Media"),
        ("数据库", "PostgreSQL 18.3，JDBC 42.7.5，HikariCP 6.2.1"),
        ("构建打包", "Maven 3.9.x，maven-shade-plugin 生成依赖完整 JAR，jpackage 生成便携 exe"),
        ("AI 与 JSON", "OkHttp 4.12.0、Jackson 2.18.2，OpenAI-compatible 接口"),
    ], [4, 11])
    add_heading(doc, "2 关键类与职责", 1)
    add_table(doc, ["类别", "类/文件", "职责"], [
        ("启动", "App、Launcher", "启动 JavaFX 应用、数据库自检、加载主窗口或降级路径"),
        ("控制器", "MainController", "目录树、缩略图、多选、搜索、右键菜单、扫描任务入口"),
        ("控制器", "DatabaseSetupDialog、SettingsController", "数据库初始化、运行配置、AI fallback 和供应商验证"),
        ("服务", "SearchService", "关键词搜索和 AI 自然语言转 SQL 搜索"),
        ("服务", "ImageServiceImpl", "删除、复制、粘贴、重命名、缩略图缓存、事务回滚"),
        ("服务", "EditService", "图片编辑版本保存、恢复和版本时间轴数据"),
        ("扫描", "DirectoryScanner、ScanTask", "扫描目录树、图片入库、AI 标签扫描"),
        ("DAO", "DatabaseConnection", "HikariCP 单例连接池、连接配置和关闭"),
        ("DAO", "TagDaoImpl、ImageDaoImpl、VersionDaoImpl", "PreparedStatement 查询、事务、标签/图片/版本持久化"),
    ], [3, 5, 7])
    add_heading(doc, "3 数据库连接与事务", 1)
    add_paragraph(doc, "系统所有主要 DAO 使用 PreparedStatement 访问数据库，参数由 JDBC 绑定，避免把用户输入直接拼接到 SQL。DatabaseConnection 使用 HikariCP 创建单例连接池，连接生命周期由 try-with-resources 关闭。")
    add_paragraph(doc, "批量重命名是事务设计的重点：应用先关闭自动提交，逐个修改文件名并更新数据库；如果中途失败，数据库 rollback，同时根据记录的路径列表把已经改动的磁盘文件名恢复。AI 标签批量入库也在事务内创建分类、标签和关联，失败时统一回滚。")
    add_picture(doc, figures["nl2sql"], "图 1 NL2SQL 安全查询链路")
    add_heading(doc, "4 数据库对象详细设计", 1)
    add_table(doc, ["类别", "对象", "详细说明"], [
        ("表", "images.thumbnail BYTEA", "缩略图缓存进入数据库，目录重复打开时优先读取，减少磁盘解码开销。"),
        ("表", "image_tags", "使用 image_id 和 tag_id 双外键实现多对多，支持 AI 与手动来源字段。"),
        ("视图", "v_image_search", "把图片、目录、标签、AI 描述聚合成面向搜索的外模式，供关键词和 NL2SQL 使用。"),
        ("索引", "idx_images_active", "对未删除图片做部分索引，匹配最常见的目录浏览查询。"),
        ("索引", "idx_tags_name_trgm / idx_ai_desc_trgm", "使用 pg_trgm GIN 支持模糊搜索，提升标签和 AI 描述检索速度。"),
        ("触发器", "trg_image_* / trg_tag_*", "自动记录图片和标签变更，补足应用层日志。"),
        ("过程", "sp_directory_report", "使用递归 CTE 统计目录子树图片数量和容量。"),
    ], [2.5, 4.5, 8])
    add_heading(doc, "5 界面与用户体验设计", 1)
    add_picture(doc, SCREENSHOTS["main"], "图 2 主界面：目录树与缩略图")
    add_picture(doc, SCREENSHOTS["db_setup"], "图 3 数据库连接与初始化向导")
    add_paragraph(doc, "主界面左侧保持完整磁盘根节点，但子目录只在展开时加载，避免启动时遍历全盘导致卡顿。数据库不可用时，标签、AI 搜索、版本历史等能力会降级，基础本地图片浏览仍可使用。")
    add_heading(doc, "6 打包与交付设计", 1)
    add_paragraph(doc, "Maven package 的主产物是 target/image-manager-1.0.0.jar，其中包含运行依赖和 sql/schema.sql、sql/data.sql。面向没有 Java 环境的电脑，scripts/package-stable.ps1 通过 jpackage 生成 DigitalImageManager.exe 和 Windows portable zip。")
    out = STAGE_DIR / f"{PRIMARY_NAME}-详细设计说明书.docx"
    doc.save(out)
    return out


def create_course_report(figures: dict[str, Path]) -> Path:
    doc = Document()
    style_document(doc, f"{PROJECT_NAME}数据库课程设计报告")
    add_heading(doc, "封面与成绩单", 1)
    add_table(doc, ["项目", "内容"], [
        ("课程设计题目", PROJECT_NAME),
        ("课程", "数据库系统基础课程设计"),
        ("小组", f"{GROUP_NAME}"),
        ("组长", "毕振岚"),
        ("报告人", "徐阳"),
        ("提交时间", DATE_TEXT),
    ], [4, 11])
    add_table(doc, ["姓名", "学号", "角色", "主要分工", "工作量比例"], [
        ("毕振岚", "202425220501", "组长", "数据库设计、需求审阅、文档统筹", "30%"),
        ("陈厚华", "202425220502", "组员", "界面交互、测试与演示支持", "25%"),
        ("徐阳", "202425220527", "项目经理/报告人", "核心实现、数据库与 AI 功能、最终交付整合", "45%"),
    ], [2.5, 3.2, 3, 5, 2])
    add_heading(doc, "1 引言", 1)
    add_heading(doc, "1.1 编写目的", 2)
    add_paragraph(doc, "本报告说明基于 PostgreSQL 的数字图像集成管理系统的需求、数据库设计、系统实现、安装使用和测试结果。报告重点证明系统不是简单的图片浏览器，而是把本地图片目录转化为可持久化、可检索、可审计、可扩展的数据库应用。")
    add_heading(doc, "1.2 定义", 2)
    add_table(doc, ["术语", "定义"], [
        ("机主", "个人电脑拥有者，也是本系统的最终用户。"),
        ("缩略图", "原图等比缩小后的预览图，本系统可用 PostgreSQL bytea 缓存。"),
        ("NL2SQL", "把自然语言搜索请求转为 PostgreSQL SELECT 查询的能力。"),
        ("fallback", "AI endpoint 逐个尝试、失败熔断和恢复的运行机制。"),
        ("逻辑删除", "图片记录保留在数据库中，但通过 is_deleted 排除出常规视图。"),
    ], [3, 12])
    add_heading(doc, "1.3 参考资料", 2)
    add_table(doc, ["资料", "用途"], [
        ("2026 数据库课程设计任务书", "确定最终提交内容、数据库设计要求和成绩构成。"),
        ("数据库课程设计报告模板", "确定报告章节结构。"),
        ("关于课程设计的重要答疑", "确定答辩时长、提交物和评分方式。"),
        ("README.md、pom.xml、schema.sql、src", "确认当前真实实现。"),
    ], [5, 10])
    add_heading(doc, "2 功能需求", 1)
    add_heading(doc, "2.1 系统范围", 2)
    add_paragraph(doc, "系统范围限定为本机桌面图片管理。它面向 Windows 用户，提供启动引导、数据库初始化、目录扫描、图片管理、AI 标签、关键词搜索、AI SQL 搜索、版本历史、幻灯片音乐和运行配置。云端/WebDAV 仅作为数据库表和依赖预留，不作为本次完整主流程。")
    add_heading(doc, "2.2 系统功能结构", 2)
    add_picture(doc, figures["function"], "图 1 系统功能结构图")
    add_heading(doc, "2.3 系统总体流程", 2)
    add_picture(doc, figures["flow"], "图 2 系统总体流程")
    add_heading(doc, "2.4 需求分析", 2)
    add_table(doc, ["功能编号", "功能名称", "说明"], REQUIREMENTS, [2.2, 4, 8.8])
    add_heading(doc, "2.5 用户界面", 2)
    add_paragraph(doc, "界面以左侧目录树、中央缩略图网格、顶部工具栏和底部状态栏为主。用户从欢迎向导进入主界面，可在数据库连接失败时继续浏览本地图片，也可以打开数据库初始化向导恢复完整能力。")
    add_picture(doc, SCREENSHOTS["welcome"], "图 3 首次启动向导")
    add_picture(doc, SCREENSHOTS["main"], "图 4 主界面")
    add_heading(doc, "2.6 数据词典", 2)
    add_table(doc, ["编号", "数据", "数据组成", "所属界面/环节", "说明"], DATA_DICTIONARY, [1.5, 2.6, 5.2, 3.2, 4.0])

    add_heading(doc, "3 数据库设计", 1)
    add_heading(doc, "3.1 PostgreSQL 选择理由", 2)
    add_paragraph(doc, "本项目选择 PostgreSQL，是因为它同时适合课程要求和系统需求。首先，系统需要完整展示主外键、唯一约束、自引用、视图、索引、触发器、存储过程和递归 CTE，PostgreSQL 对这些对象支持充分。其次，图片管理不是纯结构化数据，系统需要保存 bytea 缩略图、JSONB 编辑参数和 AI 原始 JSON，PostgreSQL 在二进制、JSON 和扩展索引方面比轻量文件存储更适合。最后，NL2SQL 搜索需要只读连接、超时控制和可被说明的视图边界，PostgreSQL 便于把自然语言查询限制在安全的 SELECT 范围内。")
    add_heading(doc, "3.2 概念结构设计", 2)
    add_picture(doc, figures["er"], "图 5 ER 与关系模式核心")
    add_paragraph(doc, "实体来源来自用户实际操作：目录树产生 directories，缩略图网格和图片信息产生 images，标签管理产生 tag_categories、tags 和 image_tags，AI 扫描产生 ai_analysis_results，搜索栏产生 search_history，编辑器时间轴产生 image_versions，删除、重命名和标签变更产生 operation_logs。")
    add_heading(doc, "3.3 逻辑结构设计", 2)
    add_table(doc, ["表名", "中文含义", "关键约束", "设计作用"], CORE_TABLES, [3, 3, 5, 5])
    add_heading(doc, "3.4 数据库模式设计", 2)
    add_table(doc, ["模式对象", "数量", "说明"], [
        ("表", "13", "覆盖目录、图片、标签、AI 结果、搜索历史、版本、操作日志、设置和扩展预留。"),
        ("索引", "19", "18 个普通索引 + 1 个唯一索引，覆盖目录过滤、哈希去重、待 AI 图片、标签模糊搜索、AI 描述搜索、当前版本和云端预留。"),
        ("视图", "4", "为常用查询、目录统计、搜索外模式和标签统计提供稳定接口。"),
        ("函数", "4", "为触发器提供日志写入逻辑。"),
        ("触发器", "5", "自动记录图片和标签变更。"),
        ("存储过程", "5", "提供月报、目录递归统计、批量改名、版本恢复和批量标签能力。"),
    ], [3, 2, 10])
    add_heading(doc, "3.5 外模式设计", 2)
    add_table(doc, ["外模式/视图", "作用"], SQL_OBJECTS["视图"], [5, 10])
    add_heading(doc, "3.6 物理结构设计", 2)
    add_paragraph(doc, "物理结构的核心是让最常用查询走索引。目录浏览使用 images(directory_id) 与 idx_images_active；待 AI 处理图片使用 idx_images_ai_pending；标签和 AI 描述模糊搜索使用 pg_trgm GIN 索引；版本时间轴使用 image_versions(image_id) 与 idx_versions_current。缩略图以 bytea 保存，避免重复打开目录时重新解码大量小图。")
    add_heading(doc, "3.7 编程性结构设计", 2)
    add_picture(doc, figures["sql_objects"], "图 6 SQL 对象与性能设计")
    for category, rows in SQL_OBJECTS.items():
        add_heading(doc, category, 3)
        add_table(doc, ["对象", "用途"], rows, [5, 10])

    add_heading(doc, "4 系统设计与实现", 1)
    add_heading(doc, "4.1 开发环境", 2)
    add_table(doc, ["类别", "选型"], [
        ("语言", "Java 21"),
        ("界面", "JavaFX 21.0.6、FXML、CSS、JavaFX Media"),
        ("数据库", "PostgreSQL 18.3"),
        ("连接", "PostgreSQL JDBC 42.7.5、HikariCP 6.2.1"),
        ("构建", "Maven 3.9.x、maven-shade-plugin、jpackage"),
        ("AI", "OpenAI-compatible Chat Completions、OkHttp、Jackson"),
    ], [4, 11])
    add_heading(doc, "4.2 整体结构与模块划分", 2)
    add_picture(doc, figures["backend"], "图 7 后台程序结构")
    add_heading(doc, "4.3 关键实现技术", 2)
    add_paragraph(doc, "目录树采用完整磁盘根节点加子目录懒加载的方式，既满足用户希望从“我的电脑”浏览的习惯，又避免启动时遍历全盘。数据库连接使用 HikariCP 连接池，所有主要 DAO 使用 PreparedStatement 绑定参数。数据库不可用时，ImageServiceImpl 可退化为直接读取磁盘图片，用户仍能完成基本浏览。")
    add_paragraph(doc, "AI 标签和 NL2SQL 是数据库部分的重点。图片经外部视觉模型分析后，原始 JSON 写入 ai_analysis_results，结构化标签写入 tag_categories、tags 和 image_tags。自然语言搜索由 AI 生成 SQL 后，系统只允许 SELECT，拒绝危险关键字，并用只读连接、5 秒超时和最大行数限制执行。")
    add_picture(doc, figures["nl2sql"], "图 8 AI 标签与 NL2SQL 安全查询链路")
    add_heading(doc, "4.4 数据库连接和事务", 2)
    add_paragraph(doc, "DatabaseConnection 负责读取数据库配置、创建 HikariDataSource 并提供 getConnection。数据库初始化向导可保存本机配置、测试连接、创建 image_manager 数据库并执行内嵌 schema.sql。批量重命名、AI 标签批量入库、AI 标签清理等流程均显式控制事务；失败时执行 rollback，避免数据库和磁盘状态不一致。")
    add_heading(doc, "4.5 系统界面设计", 2)
    add_picture(doc, SCREENSHOTS["db_setup"], "图 9 数据库连接与初始化向导")
    add_picture(doc, SCREENSHOTS["settings"], "图 10 系统设置页")
    add_picture(doc, SCREENSHOTS["editor"], "图 11 图片编辑器与版本历史")
    add_picture(doc, SCREENSHOTS["slideshow"], "图 12 幻灯片播放")

    add_heading(doc, "5 系统安装及使用说明", 1)
    add_heading(doc, "5.1 运行环境", 2)
    add_paragraph(doc, "已有 Java 环境的机器可直接运行 target/image-manager-1.0.0.jar。没有 Java 环境的 Windows 机器使用 DigitalImageManager-windows-portable.zip，解压后双击 DigitalImageManager.exe。数据库功能需要本机或局域网 PostgreSQL 服务；AI 功能需要用户自行配置 OpenAI-compatible endpoint 和密钥。")
    add_heading(doc, "5.2 配置说明", 2)
    add_table(doc, ["配置", "说明"], [
        ("数据库配置", "数据库向导写入本机外部 database.properties，不进入源码包。"),
        ("AI fallback", "设置页管理 endpoint，支持逐个测试和 last-good 恢复。"),
        ("schema.sql", "JAR 内嵌 sql/schema.sql，可由向导执行，也可手动 psql 执行。"),
        ("data.sql", "补充默认标签分类和应用设置，图片数据由扫描本机目录产生。"),
    ], [4, 11])
    add_heading(doc, "5.3 用户使用说明", 2)
    add_picture(doc, figures["demo"], "图 13 演示与使用路线")
    add_paragraph(doc, "推荐演示流程为：启动程序，进入数据库初始化向导，检测或创建数据库，选择扫描目录，查看缩略图，执行关键词搜索，切换 AI 智能搜索，打开图片编辑器保存版本，最后播放幻灯片并展示运行包。")

    add_heading(doc, "6 测试与截图", 1)
    add_table(doc, ["测试项", "验证方式", "预期结果"], [
        ("构建", "mvn -q -DskipTests compile / mvn -q -DskipTests package", "生成 image-manager-1.0.0.jar。"),
        ("单元/冒烟测试", "mvn -q test", "测试通过。"),
        ("JAR 内容", "jar tf 检查 Launcher.class、sql/schema.sql、sql/data.sql", "JAR 含主类和 SQL 资源。"),
        ("数据库初始化", "向导创建数据库并执行 schema.sql", "表、索引、视图、触发器、存储过程创建成功。"),
        ("界面", "打开欢迎、主界面、设置、数据库向导、编辑器、幻灯片截图", "无明显空白占位，图片不丢失。"),
    ], [3, 7, 5])
    add_heading(doc, "7 总结分析", 1)
    add_paragraph(doc, "本系统围绕数据库课程设计的要求，把本地图片管理中的非结构化文件整理为目录、图片、标签、AI 分析、版本和日志等结构化数据。设计中最有价值的部分不是单个界面，而是数据库对象和程序逻辑共同形成的数据闭环：启动时能建库，扫描时能入库，AI 能把图片内容转为标签，搜索能通过视图和索引返回结果，编辑能形成版本历史，触发器和日志能保留操作痕迹。")
    add_paragraph(doc, "后续如果继续扩展，可以把 cloud_sources 和 cloud_images 接入真正的 WebDAV 主流程，也可以把 image_edit_operations 的 JSONB 参数用于更精细的可逆编辑。本次最终交付保持口径克制，只把已经有代码、SQL 和截图证据支撑的能力作为主功能。")
    out = REPORT_DIR / f"{PRIMARY_NAME}-数据库课程设计报告.docx"
    doc.save(out)
    return out


def create_slide_visuals(figures: dict[str, Path]) -> list[tuple[str, Path, str]]:
    slides: list[tuple[str, Path, str]] = []
    W, H = 1920, 1080
    bg = (246, 249, 247)
    navy = (24, 50, 68)
    gray = (74, 88, 96)
    teal = (32, 128, 142)
    green = (85, 142, 121)
    orange = (186, 91, 40)

    def new_slide(title: str, subtitle: str = ""):
        img = Image.new("RGB", (W, H), bg)
        draw = ImageDraw.Draw(img)
        for i in range(0, W, 64):
            draw.line((i, 0, i, H), fill=(235, 242, 240))
        for j in range(0, H, 64):
            draw.line((0, j, W, j), fill=(235, 242, 240))
        draw.ellipse((W - 520, -220, W + 100, 420), fill=(224, 241, 238))
        draw.ellipse((-260, H - 260, 460, H + 260), fill=(249, 233, 213))
        draw.text((80, 54), title, font=pil_font(54, True), fill=navy)
        if subtitle:
            draw_wrapped(draw, (84, 128), subtitle, pil_font(24), gray, 1240, 8)
        draw.line((80, 178, 1840, 178), fill=(201, 213, 219), width=2)
        return img, draw

    def save_slide(code: str, img: Image.Image, title: str):
        path = PPT_VISUAL_DIR / f"{code}.png"
        img.save(path, quality=95)
        slides.append((title, path, code))

    img, draw = new_slide(PROJECT_NAME, "数据库系统基础课程设计答辩  |  第07组  |  JavaFX + PostgreSQL")
    draw_wrapped(draw, (90, 250), "把本地图片目录转化为可入库、可检索、可审计、可恢复的数据库应用", pil_font(48, True), navy, 760, 12)
    for i, (num, label) in enumerate([("13", "张表"), ("19", "个索引"), ("4", "个视图"), ("5", "个触发器"), ("5", "个存储过程")]):
        x = 95 + i * 150
        rounded_rectangle(draw, (x, 560, x + 130, 690), 20, (255, 255, 255), (189, 204, 212), 2)
        draw.text((x + 36, 585), num, font=pil_font(42, True), fill=orange)
        draw.text((x + 32, 638), label, font=pil_font(22), fill=gray)
    paste_cover(img, SCREENSHOTS["main"], (980, 240, 1810, 850))
    draw.text((90, 925), "组长：毕振岚    组员：陈厚华、徐阳    报告人：徐阳", font=pil_font(28), fill=gray)
    save_slide("S01", img, "题目与团队")

    img, draw = new_slide("评分点总览", "答辩和报告按教师评分表组织，数据库设计占最大权重")
    scores = [("数据库设计", 30, teal), ("功能设计", 20, green), ("后台程序设计", 15, orange), ("界面设计", 15, (120, 112, 180)), ("报告表述", 10, (115, 130, 85)), ("PPT表达", 10, (170, 110, 90))]
    total_width = 1500
    x = 160
    for label, score, color in scores:
        w = int(total_width * score / 100)
        rounded_rectangle(draw, (x, 290, x + w - 12, 460), 24, color, None, 0)
        draw.text((x + 24, 320), label, font=pil_font(28, True), fill=(255, 255, 255))
        draw.text((x + 24, 372), f"{score}分", font=pil_font(40, True), fill=(255, 255, 255))
        x += w
    notes = [
        ("数据库设计 30", "ER、关系模式、索引、视图、触发器、存储过程、递归 CTE"),
        ("功能设计 20", "启动向导、目录、图片 CRUD、AI 标签、NL2SQL、版本、幻灯片"),
        ("后台与界面 30", "HikariCP、PreparedStatement、事务回滚、fallback、真实界面截图"),
        ("表达 20", "报告按模板，PPT 少字多图，演示路线清晰"),
    ]
    for i, (t, b) in enumerate(notes):
        x = 150 + (i % 2) * 850
        y = 570 + (i // 2) * 210
        card(draw, (x, y, x + 740, y + 150), t, b, accent=scores[i][2])
    save_slide("S02", img, "评分点总览")

    img, draw = new_slide("系统功能结构", "从机主使用流程看功能，不把数据库藏在后台")
    paste_cover(img, figures["function"], (100, 230, 1820, 940))
    save_slide("S03", img, "系统功能结构")

    img, draw = new_slide("数据库设计流程", "从用户看到的数据出发，转换为实体、关系、约束和 SQL 对象")
    paste_cover(img, figures["flow"], (95, 220, 940, 885))
    paste_cover(img, figures["sql_objects"], (980, 220, 1820, 885))
    save_slide("S04", img, "数据库设计流程")

    img, draw = new_slide("ER 与关系模式", "核心是目录自引用、图片主表、标签多对多、AI 结果和版本历史")
    paste_cover(img, figures["er"], (90, 220, 1830, 930))
    save_slide("S05", img, "ER与关系模式")

    img, draw = new_slide("核心表结构与约束", "13 张表中，课程答辩重点讲清这些主表和键")
    table_cards = [
        ("directories", "parent_id 自引用\n形成目录树"),
        ("images", "directory_id 外键\nthumbnail bytea\nfile_hash / is_deleted"),
        ("image_tags", "image_id + tag_id\n多对多唯一约束"),
        ("ai_analysis_results", "image_id 唯一\nraw_response + description"),
        ("image_versions", "image_id + version_num\n版本时间轴"),
        ("operation_logs", "触发器 + 应用层\n记录增删改和标签变更"),
    ]
    for i, (t, b) in enumerate(table_cards):
        x = 110 + (i % 3) * 590
        y = 240 + (i // 3) * 300
        card(draw, (x, y, x + 500, y + 215), t, b, accent=[teal, green, orange][i % 3])
    save_slide("S06", img, "表结构与关键约束")

    img, draw = new_slide("索引、视图、触发器、存储过程", "数据库对象不是装饰，分别服务性能、外模式、审计和可编程能力")
    groups = [
        ("索引", "idx_images_active\nidx_images_ai_pending\nidx_tags_name_trgm\nidx_ai_desc_trgm\nidx_versions_current"),
        ("视图", "v_active_images\nv_directory_stats\nv_image_search\nv_tag_stats"),
        ("触发器", "trg_image_after_insert\ntrg_image_before_update\ntrg_image_after_delete\ntrg_tag_after_insert\ntrg_tag_after_delete"),
        ("存储过程", "sp_monthly_report\nsp_directory_report\nsp_batch_rename\nsp_restore_version\nsp_batch_insert_tags"),
    ]
    for i, (t, b) in enumerate(groups):
        x = 110 + (i % 2) * 870
        y = 245 + (i // 2) * 330
        card(draw, (x, y, x + 760, y + 260), t, b, accent=[teal, orange, green, (120, 112, 180)][i])
    save_slide("S07", img, "SQL对象")

    img, draw = new_slide("AI 标签与 NL2SQL 如何落库和安全查询", "AI 能力必须回到数据库：结果入库、搜索走视图、执行受安全策略约束")
    paste_cover(img, figures["nl2sql"], (110, 220, 1320, 920))
    card(draw, (1380, 260, 1815, 420), "落库", "ai_analysis_results 保存原始 JSON 与描述；image_tags 保存标签关联", accent=teal)
    card(draw, (1380, 470, 1815, 630), "查询", "优先使用 v_image_search，图片 id 必须作为第一列", accent=green)
    card(draw, (1380, 680, 1815, 840), "安全", "只允许 SELECT，禁危险词，只读连接，5 秒超时，最多 1000 行", accent=orange)
    save_slide("S08", img, "AI与NL2SQL")

    img, draw = new_slide("后台程序设计", "分层、连接池、事务和 fallback 保证系统可运行、可解释、可答问")
    paste_cover(img, figures["backend"], (100, 225, 1080, 915))
    paste_cover(img, SCREENSHOTS["settings"], (1140, 225, 1815, 570))
    paste_cover(img, SCREENSHOTS["db_setup"], (1140, 610, 1815, 915))
    save_slide("S09", img, "后台程序设计")

    img, draw = new_slide("核心界面截图", "界面评分看真实可用状态：开屏、主界面、数据库向导、编辑、幻灯片")
    shots = [SCREENSHOTS["welcome"], SCREENSHOTS["main"], SCREENSHOTS["db_setup"], SCREENSHOTS["editor"], SCREENSHOTS["slideshow"], SCREENSHOTS["rename"]]
    labels = ["首次启动", "主界面", "数据库向导", "编辑与版本", "幻灯片", "批量重命名"]
    for i, (shot, label) in enumerate(zip(shots, labels)):
        x = 100 + (i % 3) * 600
        y = 230 + (i // 3) * 340
        paste_cover(img, shot, (x, y, x + 520, y + 270))
        draw.text((x + 12, y + 282), label, font=pil_font(22, True), fill=navy)
    save_slide("S10", img, "核心界面截图")

    img, draw = new_slide("代码与数据库演示路线", "先跑可交付物，再证明数据库对象和功能闭环")
    paste_cover(img, figures["demo"], (100, 220, 1280, 920))
    commands = "java -jar app.jar\nOpen DigitalImageManager.exe\n\npsql checklist:\n\\dt\n\\dv\n\\df\n\\d+ v_image_search"
    rounded_rectangle(draw, (1350, 260, 1815, 610), 24, (28, 42, 52), None, 0)
    draw_wrapped(draw, (1385, 295), commands, pil_font(23, mono=True), (232, 246, 237), 390, 10)
    card(draw, (1350, 675, 1815, 845), "现场重点", "初始化数据库 -> 扫描目录 -> 搜索 -> 编辑版本 -> 查看交付包", accent=orange)
    save_slide("S11", img, "演示路线")

    img, draw = new_slide("创新点与总结", "本项目的分数点集中在数据库含金量和可运行交付")
    points = [
        ("开屏引导", "第一次启动即选择目录，提示 AI 成本和数据库状态"),
        ("离线降级", "数据库不可用仍能浏览本地图片，降低演示风险"),
        ("逐个供应商验证", "AI fallback 支持测试、熔断和 last-good 恢复"),
        ("AI 转 SQL 搜索", "自然语言落到 v_image_search，只读安全执行"),
        ("数据库建立引导", "向导一键创建 image_manager 并执行内嵌 schema.sql"),
        ("交付完整", "JAR、portable zip/exe、源码包、验证日志齐全"),
    ]
    for i, (t, b) in enumerate(points):
        x = 115 + (i % 3) * 590
        y = 245 + (i // 3) * 300
        card(draw, (x, y, x + 505, y + 210), t, b, accent=[teal, green, orange][i % 3])
    draw_wrapped(draw, (120, 925), "边界说明：WebDAV/云端只作为扩展预留，不作为本次完整主流程。", pil_font(28, True), navy, 1500, 8)
    save_slide("S12", img, "创新点与总结")

    return slides


def create_ppt(slides: list[tuple[str, Path, str]]) -> Path:
    prs = Presentation()
    prs.slide_width = PptInches(13.333333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]
    for title, path, code in slides:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(path), 0, 0, width=prs.slide_width, height=prs.slide_height)
        notes = slide.notes_slide.notes_text_frame
        notes.text = f"{code} {title}"
    out = PPT_DIR / f"{PRIMARY_NAME}-数据库课程设计答辩PPT.pptx"
    prs.save(out)
    return out


def create_ppt_planning(slides: list[tuple[str, Path, str]]) -> None:
    content_report = PPT_DIR / "content_report.md"
    content_report.write_text(
        f"""# content_report

## I. Source Status

source_level: complete_report

本 PPT 的内容基础来自教师任务书、课程报告模板、重要答疑、打分依据图片、需求分析、概要设计、详细设计、README、schema.sql、src 和现有真实截图。当前材料足以直接生成答辩用 PPT，不再扩写无证据的功能。

## II. Content Thesis

这套答辩 PPT 要证明：{PROJECT_NAME} 不是普通图片浏览器，而是以 PostgreSQL 为核心的桌面数据库应用。系统把目录、图片元数据、缩略图、AI 标签、自然语言搜索、版本历史和操作日志组织为可解释的数据模型，并提供可运行的 JAR 与 Windows portable 包。

## III. Narrative Body

PPT 先说明评分点，再展示系统功能结构和数据库设计流程。中段集中讲 ER、表结构、索引、视图、触发器、存储过程、AI 标签落库和 NL2SQL 安全查询。后段用真实界面截图证明系统可用，并给出现场演示路线。云端和 WebDAV 只作为扩展预留，不放入主功能承诺。

## IV. Section Candidates

1. 题目与评分点。
2. 功能结构和数据库设计流程。
3. ER、表、SQL 对象和 AI/NL2SQL。
4. 后台程序、界面截图和演示路线。
5. 创新点与总结。

## V. Visualizable Content

可视化内容包括功能结构图、ER 图、SQL 对象总览、NL2SQL 链路图、后台架构图、演示路线图和真实软件截图。
""",
        encoding="utf-8",
    )
    (PPT_DIR / "design_spec.md").write_text(
        """# design_spec

## I. Project Information

项目名：基于 PostgreSQL 的数字图像集成管理系统。用途：数据库系统基础课程设计答辩。观众：教师和同学。时长：不超过 15 分钟。比例：16:9。

## II. Narrative Spine

叙事主线是“评分点 -> 数据库含金量 -> 可运行系统 -> 演示路线”。数据库设计是核心，功能和界面服务于数据库应用证明。

## III. Style Direction

采用常规专业路线，明亮克制、图像优先、少字、结构化信息容器。页面使用真实截图和生成的数据库图表，避免营销式大段文字。

## IV. Global Design Principles

背景明亮，主色使用深墨色、青绿色、橙色和浅灰；正文页保持高信息密度但不堆满文字；每页只保留一个主消息。

## V. Constraints

不写云端/WebDAV 为已完成主流程；不写 AI 为本地内置模型；不写无证据的固定上限。所有可见功能要能在 README、schema.sql 或源码中找到证据。
""",
        encoding="utf-8",
    )
    lines = ["# slide_blueprint\n"]
    for title, path, code in slides:
        lines.append(f"## {code}\n")
        lines.append(f"- page_role: 答辩页\n- title: {title}\n- core_message: 围绕数据库课程设计评分点展示 {title}。\n- content_basis_binding: 教师任务书、当前 README、schema.sql、src 和真实截图。\n- claim_status: user_provided_and_code_verified\n- page_rhythm: 图像优先，少字说明。\n- text_visual_balance: 以整页图像为主。\n- visual_strategy: 使用生成图表或真实界面截图。\n- continuity_inheritance: 明亮背景、青绿/橙色强调、结构化卡片。\n- preview_evidence_binding: 采用本地生成的 16:9 页面视觉 `{path.name}`。\n")
    (PPT_DIR / "slide_blueprint.md").write_text("\n".join(lines), encoding="utf-8")
    (PPT_DIR / "spec_lock.md").write_text(
        """# spec_lock

## canvas

- format: PPTX
- ratio: 16:9
- page_visuals: one final image per slide

## visual_system

- brightness_world: bright professional academic deck
- palette_roles: off-white base, deep ink title, teal/green/orange accents, light gray grid
- container_grammar: rounded white cards and full-slide image panels
- emphasis_tone: database-first, precise, defendable

## content_grounding

- allowed: facts supported by README, schema.sql, src, screenshots, teacher materials
- forbidden: cloud/WebDAV as completed main flow; AI as local model; unsupported limits

## generation_metadata

Slide IDs and filenames stay in planning files and image filenames, not in visible slide copy except normal page order.
""",
        encoding="utf-8",
    )


def create_speech_doc(slides: list[tuple[str, Path, str]], figures: dict[str, Path]) -> Path:
    doc = Document()
    style_document(doc, f"{PROJECT_NAME}答辩打印讲稿")
    add_heading(doc, "使用方式", 1)
    add_paragraph(doc, "答辩顺序建议为先讲 PPT，再运行 JAR 或 Windows portable exe 演示。PPT 讲述按下表压到约 10 分 15 秒，现场演示控制在 3 到 4 分钟，总时长留出提问余量。讲稿不是逐字背诵稿，而是手持提醒稿，重点帮助回答数据库老师可能追问的底层问题。")
    add_heading(doc, "每页讲稿", 1)
    page_notes = [
        ("S01", "题目与团队", "30 秒", "说明系统是 JavaFX + PostgreSQL 桌面图片管理系统，不是普通图片浏览器；先报成员和报告人。"),
        ("S02", "评分点总览", "35 秒", "直接对齐评分表，强调数据库设计 30 分是主线，后面所有页都围绕这个主线展开。"),
        ("S03", "系统功能结构", "40 秒", "只按主流程带过：启动、目录、图片操作、AI 标签、搜索、版本、幻灯片和设置。"),
        ("S04", "数据库设计流程", "50 秒", "解释如何从用户看到的数据抽象实体，再落到表、约束、视图、索引、触发器和过程。"),
        ("S05", "ER 与关系模式", "70 秒", "重点讲 directories 自引用、images 主表、标签多对多、AI 结果一对一、版本一对多。"),
        ("S06", "核心表结构与约束", "70 秒", "把主键、外键、唯一约束、bytea、逻辑删除讲清楚。"),
        ("S07", "SQL 对象", "70 秒", "说明索引提升查询，视图作为外模式，触发器记录审计，存储过程展示数据库可编程能力。"),
        ("S08", "AI 与 NL2SQL", "75 秒", "说明 AI 不是随便查数据库，而是结果先入库，查询必须 SELECT、只读、超时、限行。"),
        ("S09", "后台程序设计", "55 秒", "讲 HikariCP、PreparedStatement、事务回滚、后台任务和 fallback。"),
        ("S10", "核心界面截图", "45 秒", "快速证明界面可用：欢迎、主界面、数据库向导、编辑器、幻灯片、批量重命名。"),
        ("S11", "演示路线", "35 秒", "把接下来 4 个动作说清：初始化、搜索、版本、交付。"),
        ("S12", "创新点与总结", "40 秒", "收束到开屏引导、离线降级、逐个供应商验证、AI 转 SQL、数据库建立引导。"),
    ]
    add_table(doc, ["页码", "页面", "时间", "讲述重点"], page_notes, [1.5, 3, 2, 9])
    add_heading(doc, "演示点击路线", 1)
    add_table(doc, ["步骤", "操作", "说明"], [
        ("1", "运行 `java -jar target/image-manager-1.0.0.jar` 或双击 `DigitalImageManager.exe`，进入数据库向导", "用 1 分钟证明运行包可启动，并说明向导执行内嵌 `schema.sql` 和 `data.sql`。"),
        ("2", "选择扫描目录，展示主界面并执行关键词/NL2SQL 搜索", "用 1 分钟讲完整磁盘树懒加载、缩略图入库、`v_image_search`、只读 SELECT 和 `search_history`。"),
        ("3", "打开图片编辑器保存版本，再展示版本时间轴", "用 1 分钟讲 `image_versions`、当前版本标记和 `.versions` 文件夹。"),
        ("4", "展示最终交付目录和验证日志", "用 30 到 60 秒证明 JAR、portable zip/exe、源码包、报告、PPT、讲稿、验证日志齐全。"),
    ], [1.5, 6, 7.5])
    add_heading(doc, "数据库问答准备", 1)
    qa = [
        ("为什么使用 PostgreSQL？", "需要完整展示主外键、视图、索引、触发器、存储过程和递归 CTE；同时支持 bytea 和 JSONB，适合缩略图、AI JSON 和编辑参数。"),
        ("directories 的自引用怎么工作？", "`parent_id` 引用 `directories(id)`，根目录为 NULL，子目录通过 parent_id 形成树。递归 CTE 可统计某目录下所有子目录图片。"),
        ("images 表为什么有 bytea？", "缩略图是高频读取的小二进制对象，缓存到 `thumbnail BYTEA` 后，目录重复打开时可以优先读数据库，减少磁盘解码。"),
        ("标签为什么拆成三张表？", "`tag_categories` 管分类，`tags` 管具体词，`image_tags` 管图片与标签的多对多关系，避免同一个标签重复存到每张图片记录里。"),
        ("索引为什么这样建？", "目录浏览最频繁，所以建 `images(directory_id)` 和 `idx_images_active`；AI 扫描要快速找未处理图片，所以建 `idx_images_ai_pending`；标签和描述要模糊搜，所以用 `idx_tags_name_trgm` 与 `idx_ai_desc_trgm`；版本时间轴要找当前版本，所以建 `idx_versions_current`；`idx_images_hash` 是唯一部分索引，用于哈希去重。"),
        ("NL2SQL 怎样防危险 SQL？", "先要求 AI 只生成 SELECT，再在 Java 侧检查首词和危险关键字；执行时设置只读连接、5 秒超时和 1000 行上限。"),
        ("触发器记录什么？", "图片新增、改名、删除状态变化和标签关联变化都会写入 `operation_logs`，用于审计。"),
        ("存储过程是不是都被前台调用？", "不夸大。部分前台主流程由 Java Service/DAO 事务完成；存储过程体现数据库层可编程能力和扩展接口，例如目录报表、版本恢复、批量标签。"),
        ("数据库不可用为什么还能打开？", "基础浏览可直接读取磁盘文件；数据库相关的标签、AI、搜索历史和版本功能降级，避免演示时被数据库连接阻断。"),
        ("WebDAV 完成了吗？", "没有作为主流程完成。当前有依赖和 cloud_sources/cloud_images 表，是后续扩展预留，报告和 PPT 中只按预留说明。"),
    ]
    add_table(doc, ["问题", "回答要点"], qa, [5, 10])
    add_heading(doc, "表结构速记图", 1)
    add_picture(doc, figures["er"], "图 1 ER 与关系模式速记")
    out = SPEECH_DIR / f"{PRIMARY_NAME}-数据库课程设计答辩讲稿.docx"
    doc.save(out)
    return out


def create_evidence_docs(outputs: dict[str, str], figures: dict[str, Path]) -> None:
    (EVIDENCE_DIR / "打分依据识图转写.md").write_text(
        """# 打分依据识图转写

来源：`docs/数据库系统基础/2026春-数据库系统课程设计教师提供的资料/txt中提到的打分依据.png`

| 项目 | 分值 |
|---|---:|
| 数据库设计 | 最高 30 分 |
| 后台程序设计 | 最高 15 分 |
| 功能设计 | 最高 20 分 |
| 界面设计 | 最高 15 分 |
| 报告表述水平 | 最高 10 分 |
| PPT 表达能力 | 最高 10 分 |
| 总分 | 100 分 |

图片中另有字段：课程设计题目、组长姓名、报告人姓名、总分、简单评价。
""",
        encoding="utf-8",
    )
    rows = "\n".join(f"| {name} | `{path}` |" for name, path in outputs.items())
    (EVIDENCE_DIR / "最终交付清单.md").write_text(
        f"""# 最终交付清单

## 正式产物

| 产物 | 路径 |
|---|---|
{rows}

## 目录要求

| 目录 | 内容 | 状态 |
|---|---|---|
| `01_阶段文档` | 需求分析、概要设计、详细设计 DOCX 交付版 | 已生成 |
| `02_课程报告` | 按课程报告模板生成的 DOCX | 已生成 |
| `03_答辩PPT` | 16:9 答辩 PPT、页面视觉图和规划文件 | 已生成 |
| `04_打印讲稿` | 适合打印手持的答辩稿 DOCX | 已生成 |
| `05_源码与运行包` | 源码包、JAR、portable zip/exe | 由最终验证脚本同步并校验 |
| `06_证据与清单` | 清单、图表、抽取材料、验证日志 | 由最终验证脚本写入验证结论 |

## 口径边界

1. WebDAV/云端只作为扩展预留。
2. AI 需要外部 OpenAI-compatible endpoint 和用户配置，不写成本地内置模型。
3. 存储过程作为数据库课程对象和扩展能力，不写成所有前台主流程都直接调用。
""",
        encoding="utf-8",
    )
    (EVIDENCE_DIR / "写作依据与取舍说明.md").write_text(
        """# 写作依据与取舍说明

正式文档依据教师任务书、课程报告模板、重要答疑、打分依据图片、需求分析、概要设计、详细设计、README、pom.xml、schema.sql、src 和 package-stable.ps1 生成。过程记录中的具体原话不进入正式材料，只转化为“数据库价值必须讲清楚”“AI 标签、SQL 检索、索引和版本历史需要体现数据库含金量”等设计依据。

剔除的旧口径包括：JavaFX 26、SQLite 已实现、AI Key 入库、完整 WebDAV 云端主流程、全网全盘搜索、语音搜索、AI 标签置信度由模型返回、每张图片强制最多 50 个版本、搜索历史固定保留最近 500 条。
""",
        encoding="utf-8",
    )


def create_source_zip() -> Path:
    out = PACKAGE_DIR / f"{PRIMARY_NAME}-数据库课程设计源代码.zip"
    include_roots = ["src", "sql", "assets", "scripts"]
    include_files = ["README.md", "pom.xml", ".gitignore", ".gitattributes", "WORK_LOG.md", "AGENTS.md"]
    excluded_parts = {
        ".git", "target", "logs", ".idea", ".vscode", "node_modules", "__pycache__",
        "旧的最终交付，现在不使用", "收齐其他小组作业",
    }
    with zipfile.ZipFile(out, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for file in include_files:
            path = ROOT / file
            if path.exists():
                zf.write(path, file)
        for root_name in include_roots:
            root_path = ROOT / root_name
            if not root_path.exists():
                continue
            for path in root_path.rglob("*"):
                if not path.is_file():
                    continue
                rel = path.relative_to(ROOT)
                if any(part in excluded_parts for part in rel.parts):
                    continue
                if any(str(rel).lower().endswith(ext) for ext in [".key", ".pem", ".env"]):
                    continue
                zf.write(path, rel.as_posix())
    return out


def main() -> None:
    ensure_dirs()
    figures = generate_figures()
    outputs: dict[str, str] = {}
    req = create_stage_requirement(figures)
    outputs["系统需求分析说明书"] = str(req.relative_to(ROOT))
    overview = create_stage_overview(figures)
    outputs["概要设计说明书"] = str(overview.relative_to(ROOT))
    detail = create_stage_detail(figures)
    outputs["详细设计说明书"] = str(detail.relative_to(ROOT))
    report = create_course_report(figures)
    outputs["数据库课程设计报告"] = str(report.relative_to(ROOT))
    slides = create_slide_visuals(figures)
    create_ppt_planning(slides)
    ppt = create_ppt(slides)
    outputs["数据库课程设计答辩PPT"] = str(ppt.relative_to(ROOT))
    speech = create_speech_doc(slides, figures)
    outputs["答辩打印讲稿"] = str(speech.relative_to(ROOT))
    source_zip = create_source_zip()
    outputs["源码快照"] = str(source_zip.relative_to(ROOT))
    create_evidence_docs(outputs, figures)
    print(json.dumps(outputs, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
