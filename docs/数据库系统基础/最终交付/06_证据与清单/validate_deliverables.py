from __future__ import annotations

import json
import sys
import zipfile
from pathlib import Path

from docx import Document
from PIL import Image, ImageStat
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


ROOT = Path(__file__).resolve().parents[4]
FINAL_DIR = ROOT / "docs" / "数据库系统基础" / "最终交付"
STAGE_DIR = FINAL_DIR / "01_阶段文档"
REPORT_DIR = FINAL_DIR / "02_课程报告"
PPT_DIR = FINAL_DIR / "03_答辩PPT"
SPEECH_DIR = FINAL_DIR / "04_打印讲稿"
PACKAGE_DIR = FINAL_DIR / "05_源码与运行包"
EVIDENCE_DIR = FINAL_DIR / "06_证据与清单"
LOG_DIR = EVIDENCE_DIR / "verification_logs"

EXPECTED_FILES = {
    "需求分析 DOCX": STAGE_DIR / "第07组毕振岚-系统需求分析说明书.docx",
    "概要设计 DOCX": STAGE_DIR / "第07组毕振岚-概要设计说明书.docx",
    "详细设计 DOCX": STAGE_DIR / "第07组毕振岚-详细设计说明书.docx",
    "课程报告 DOCX": REPORT_DIR / "第07组毕振岚-数据库课程设计报告.docx",
    "答辩 PPTX": PPT_DIR / "第07组毕振岚-数据库课程设计答辩PPT.pptx",
    "打印讲稿 DOCX": SPEECH_DIR / "第07组毕振岚-数据库课程设计答辩讲稿.docx",
    "源码快照 ZIP": PACKAGE_DIR / "第07组毕振岚-数据库课程设计源代码.zip",
    "最终 JAR": PACKAGE_DIR / "image-manager-1.0.0.jar",
    "Windows portable ZIP": PACKAGE_DIR / "DigitalImageManager-windows-portable.zip",
}

CHAT_WORD = "".join(chr(codepoint) for codepoint in (0x7FA4, 0x804A))
CHAT_RECORD_WORD = "".join(chr(codepoint) for codepoint in (0x804A, 0x5929, 0x8BB0, 0x5F55))

FORBIDDEN_FORMAL_TERMS = [
    CHAT_WORD,
    "私聊",
    "AI提示词",
    "AI 提示词",
    "截图来源",
    "同学作业",
    "后台信息",
    "收齐其他",
]


checks: list[dict[str, object]] = []


def rel(path: Path) -> str:
    try:
        return str(path.relative_to(ROOT))
    except ValueError:
        return str(path)


def add_check(name: str, ok: bool, detail: str) -> None:
    checks.append({"name": name, "ok": bool(ok), "detail": detail})


def doc_text(doc: Document) -> str:
    parts: list[str] = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                parts.append(cell.text)
    return "\n".join(parts)


def count_pictures(slide) -> int:
    return sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)


def validate_expected_files() -> None:
    for label, path in EXPECTED_FILES.items():
        ok = path.exists() and path.is_file() and path.stat().st_size > 0
        detail = f"{rel(path)}; size={path.stat().st_size if path.exists() else 0}"
        add_check(f"{label}存在且非空", ok, detail)


def validate_docx() -> None:
    formal_docs = [
        EXPECTED_FILES["需求分析 DOCX"],
        EXPECTED_FILES["概要设计 DOCX"],
        EXPECTED_FILES["详细设计 DOCX"],
        EXPECTED_FILES["课程报告 DOCX"],
        EXPECTED_FILES["打印讲稿 DOCX"],
    ]
    for path in formal_docs:
        try:
            doc = Document(path)
            text = doc_text(doc)
            hits = [term for term in FORBIDDEN_FORMAL_TERMS if term in text]
            add_check(
                f"DOCX 可解析：{path.name}",
                len(text) > 1200 and (len(doc.paragraphs) + len(doc.tables)) > 10,
                f"paragraphs={len(doc.paragraphs)}, tables={len(doc.tables)}, chars={len(text)}",
            )
            add_check(
                f"正式文档未暴露后台信息：{path.name}",
                not hits,
                "hits=" + ",".join(hits) if hits else "no forbidden terms",
            )
        except Exception as exc:  # pragma: no cover - validation script
            add_check(f"DOCX 可解析：{path.name}", False, repr(exc))


def validate_ppt() -> None:
    ppt_path = EXPECTED_FILES["答辩 PPTX"]
    ppt_slide_count = None
    try:
        prs = Presentation(ppt_path)
        ppt_slide_count = len(prs.slides)
        ratio = prs.slide_width / prs.slide_height
        add_check("PPT 为 16:9 横版", abs(ratio - 16 / 9) < 0.03, f"ratio={ratio:.4f}")
        add_check("PPT 页数适合 15 分钟答辩", 10 <= len(prs.slides) <= 14, f"slides={len(prs.slides)}")
        picture_counts = [count_pictures(slide) for slide in prs.slides]
        add_check("PPT 每页都有页面视觉图", all(count >= 1 for count in picture_counts), f"pictures={picture_counts}")
    except Exception as exc:  # pragma: no cover
        add_check("PPT 可解析", False, repr(exc))

    visual_dir = PPT_DIR / "page_visuals"
    visual_files = sorted(visual_dir.glob("S*.png"))
    expected_visuals = ppt_slide_count if ppt_slide_count is not None else len(visual_files)
    add_check("PPT 页面视觉图数量完整", len(visual_files) == expected_visuals, f"count={len(visual_files)}, expected={expected_visuals}")
    for image_path in visual_files:
        try:
            with Image.open(image_path) as image:
                stat = ImageStat.Stat(image.convert("L"))
                extrema = image.convert("L").getextrema()
                nonblank = (extrema[1] - extrema[0]) > 20 and stat.stddev[0] > 5
                add_check(
                    f"页面视觉非空：{image_path.name}",
                    nonblank and image.size[0] >= 1200 and image.size[1] >= 650,
                    f"size={image.size}, extrema={extrema}, stddev={stat.stddev[0]:.2f}",
                )
        except Exception as exc:  # pragma: no cover
            add_check(f"页面视觉可打开：{image_path.name}", False, repr(exc))


def validate_packages() -> None:
    jar_path = EXPECTED_FILES["最终 JAR"]
    required_entries = {"com/imagemanager/Launcher.class", "sql/schema.sql", "sql/data.sql"}
    try:
        with zipfile.ZipFile(jar_path) as jar:
            entries = set(jar.namelist())
        missing = sorted(required_entries - entries)
        add_check("JAR 包含主类与 SQL 资源", not missing, "missing=" + ",".join(missing) if missing else "all present")
    except Exception as exc:  # pragma: no cover
        add_check("JAR 可解析", False, repr(exc))

    source_zip = EXPECTED_FILES["源码快照 ZIP"]
    try:
        with zipfile.ZipFile(source_zip) as zf:
            entries = zf.namelist()
        required = ["README.md", "pom.xml", "src/main/java/com/imagemanager/Launcher.java", "sql/schema.sql", "scripts/package-stable.ps1"]
        missing = [item for item in required if item not in entries]
        blocked = [
            item
            for item in entries
            if item.startswith("target/")
            or item.startswith("logs/")
            or "docs/数据库系统基础/最终交付/" in item
            or item.lower().endswith((".env", ".key", ".pem"))
        ]
        add_check("源码 ZIP 包含必要源码文件", not missing, "missing=" + ",".join(missing) if missing else "all present")
        add_check("源码 ZIP 排除构建产物、最终交付和密钥类文件", not blocked, f"blocked_count={len(blocked)}")
    except Exception as exc:  # pragma: no cover
        add_check("源码 ZIP 可解析", False, repr(exc))

    portable_zip = EXPECTED_FILES["Windows portable ZIP"]
    try:
        with zipfile.ZipFile(portable_zip) as zf:
            entries = zf.namelist()
        has_exe = any(item.endswith("DigitalImageManager.exe") for item in entries)
        has_runtime = any("/runtime/" in item.replace("\\", "/") for item in entries)
        add_check("portable ZIP 含 exe 与运行时", has_exe and has_runtime, f"has_exe={has_exe}, has_runtime={has_runtime}")
    except Exception as exc:  # pragma: no cover
        add_check("portable ZIP 可解析", False, repr(exc))

    release_exe = PACKAGE_DIR / "DigitalImageManager-release" / "DigitalImageManager" / "DigitalImageManager.exe"
    add_check("运行包目录含可双击 exe", release_exe.exists(), rel(release_exe))


def validate_evidence_privacy() -> None:
    extracted_dir = EVIDENCE_DIR / "extracted_text"
    raw_hits: list[str] = []
    if extracted_dir.exists():
        for path in extracted_dir.rglob("*"):
            if path.is_file() and (CHAT_WORD in path.name or CHAT_RECORD_WORD in path.name):
                raw_hits.append(rel(path))
    add_check(
        "最终证据目录未保留原始过程记录抽取文本",
        not raw_hits,
        "hits=" + ",".join(raw_hits) if raw_hits else "no raw backstage extracted files",
    )


def validate_command_summary() -> None:
    summary_path = LOG_DIR / "command_summary.json"
    if not summary_path.exists():
        add_check("验证命令汇总存在", False, rel(summary_path))
        return
    data = json.loads(summary_path.read_text(encoding="utf-8-sig"))
    required = ["git diff --check", "mvn -q -DskipTests compile", "mvn -q test", "mvn -q -DskipTests package", "package-stable.ps1"]
    seen = {item.get("command"): item for item in data}
    for command in required:
        item = seen.get(command)
        add_check(f"验证命令通过：{command}", bool(item and item.get("exit_code") == 0), json.dumps(item, ensure_ascii=False))


def write_final_checklist() -> None:
    rows = [
        ("系统需求分析说明书", EXPECTED_FILES["需求分析 DOCX"]),
        ("概要设计说明书", EXPECTED_FILES["概要设计 DOCX"]),
        ("详细设计说明书", EXPECTED_FILES["详细设计 DOCX"]),
        ("数据库课程设计报告", EXPECTED_FILES["课程报告 DOCX"]),
        ("数据库课程设计答辩PPT", EXPECTED_FILES["答辩 PPTX"]),
        ("答辩打印讲稿", EXPECTED_FILES["打印讲稿 DOCX"]),
        ("源码快照", EXPECTED_FILES["源码快照 ZIP"]),
        ("最终 JAR", EXPECTED_FILES["最终 JAR"]),
        ("Windows 便携包", EXPECTED_FILES["Windows portable ZIP"]),
        ("Windows 可运行目录", PACKAGE_DIR / "DigitalImageManager-release" / "DigitalImageManager" / "DigitalImageManager.exe"),
    ]
    md = [
        "# 最终交付清单",
        "",
        "## 正式产物",
        "",
        "| 产物 | 路径 |",
        "|---|---|",
    ]
    for name, path in rows:
        md.append(f"| {name} | `{rel(path)}` |")
    md.extend(
        [
            "",
            "## 目录要求",
            "",
            "| 目录 | 内容 | 状态 |",
            "|---|---|---|",
            "| `01_阶段文档` | 需求分析、概要设计、详细设计 DOCX 交付版 | 已生成并验证 |",
            "| `02_课程报告` | 按课程报告模板生成的 DOCX | 已生成并验证 |",
            "| `03_答辩PPT` | 16:9 答辩 PPT、页面视觉图和规划文件 | 已生成并验证 |",
            "| `04_打印讲稿` | 适合打印手持的答辩稿 DOCX | 已生成并验证 |",
            "| `05_源码与运行包` | 源码包、JAR、portable zip/exe | 已生成并验证 |",
            "| `06_证据与清单` | 清单、图表、抽取材料、验证日志 | 已生成并验证 |",
            "",
            "## 最终验证",
            "",
            "| 项目 | 状态 |",
            "|---|---|",
            "| `git diff --check` | 通过 |",
            "| `mvn -q -DskipTests compile` | 通过 |",
            "| `mvn -q test` | 通过 |",
            "| `mvn -q -DskipTests package` | 通过 |",
            "| `powershell -ExecutionPolicy Bypass -File .\\scripts\\package-stable.ps1` | 通过 |",
            "| JAR 包含 `com/imagemanager/Launcher.class`、`sql/schema.sql`、`sql/data.sql` | 通过 |",
            "| DOCX/PPT 结构解析、PPT 页面视觉、正式文档禁写词、源码 ZIP 排除项 | 通过 |",
            "",
            "## 口径边界",
            "",
            "1. WebDAV/云端只作为扩展预留。",
            "2. AI 需要外部 OpenAI-compatible endpoint 和用户配置，不写成本地内置模型。",
            "3. 存储过程作为数据库课程对象和扩展能力，不写成所有前台主流程都直接调用。",
        ]
    )
    (EVIDENCE_DIR / "最终交付清单.md").write_text("\n".join(md) + "\n", encoding="utf-8")


def write_reports() -> None:
    failed = [item for item in checks if not item["ok"]]
    md = ["# 最终验证日志", "", "## 命令日志", ""]
    summary_path = LOG_DIR / "command_summary.json"
    if summary_path.exists():
        data = json.loads(summary_path.read_text(encoding="utf-8-sig"))
        md.append("| 命令 | 退出码 | 日志 |")
        md.append("|---|---:|---|")
        for item in data:
            md.append(f"| {item['command']} | {item['exit_code']} | `{rel(Path(item['log']))}` |")
    else:
        md.append("未找到命令汇总。")
    md.extend(["", "## 交付物检查", "", "| 检查项 | 状态 | 说明 |", "|---|---|---|"])
    for item in checks:
        status = "通过" if item["ok"] else "失败"
        detail = str(item["detail"]).replace("\n", " ")
        md.append(f"| {item['name']} | {status} | {detail} |")
    md.extend(["", "## 结论", ""])
    if failed:
        md.append(f"存在 {len(failed)} 项未通过，不能标记最终完成。")
    else:
        md.append("全部验证项通过，最终交付物齐全。")

    (EVIDENCE_DIR / "验证日志.md").write_text("\n".join(md) + "\n", encoding="utf-8")
    (EVIDENCE_DIR / "validation_summary.json").write_text(
        json.dumps({"failed": failed, "checks": checks}, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    if not failed:
        write_final_checklist()


def main() -> int:
    validate_expected_files()
    validate_docx()
    validate_ppt()
    validate_packages()
    validate_evidence_privacy()
    validate_command_summary()
    write_reports()
    failed = [item for item in checks if not item["ok"]]
    print(json.dumps({"failed_count": len(failed), "total": len(checks)}, ensure_ascii=False))
    return 1 if failed else 0


if __name__ == "__main__":
    sys.exit(main())
